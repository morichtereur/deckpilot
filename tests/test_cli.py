"""End-to-end: the commands a reader of the README will actually run."""

import logging

import pytest
from click.testing import CliRunner
from pptx import Presentation

from deckpilot.cli import cli
from deckpilot.data.generate import build_programme
from deckpilot.renderer.qa import check_deck
from deckpilot.specgen.samples import SAMPLES


@pytest.fixture
def runner() -> CliRunner:
    return CliRunner()


def test_demo_builds_a_deck_without_an_api_key(runner, tmp_path, monkeypatch):
    monkeypatch.delenv("ANTHROPIC_API_KEY", raising=False)
    out = tmp_path / "deck.pptx"
    result = runner.invoke(cli, ["demo", "--out", str(out)])
    assert result.exit_code == 0, result.output
    assert out.exists()

    prs = Presentation(str(out))
    assert len(prs.slides) >= 8
    assert "Geometry check: clean." in result.output


def test_the_demo_deck_has_no_overflow_warnings(runner, tmp_path, caplog):
    out = tmp_path / "deck.pptx"
    with caplog.at_level(logging.WARNING, logger="deckpilot.renderer"):
        result = runner.invoke(cli, ["demo", "--out", str(out)])
    assert result.exit_code == 0
    assert "truncated" not in caplog.text, caplog.text


def test_the_demo_deck_is_geometrically_clean(runner, tmp_path):
    out = tmp_path / "deck.pptx"
    assert runner.invoke(cli, ["demo", "--out", str(out)]).exit_code == 0
    assert [f for f in check_deck(out) if f.severity == "error"] == []


def test_build_reads_a_programme_file(runner, tmp_path):
    data = tmp_path / "programme.json"
    build_programme().save(data)
    out = tmp_path / "deck.pptx"
    spec_out = tmp_path / "spec.json"
    result = runner.invoke(
        cli,
        ["build", "--data", str(data), "--week", "2026-W34", "--out", str(out),
         "--spec-out", str(spec_out)],
    )
    assert result.exit_code == 0, result.output
    assert out.exists() and spec_out.exists()


def test_build_refuses_a_week_it_has_no_data_for(runner, tmp_path):
    data = tmp_path / "programme.json"
    build_programme().save(data)
    result = runner.invoke(
        cli, ["build", "--data", str(data), "--week", "1999-W01",
              "--out", str(tmp_path / "d.pptx")]
    )
    assert result.exit_code != 0
    assert "no status reported" in str(result.exception)


@pytest.mark.parametrize("layout", sorted(SAMPLES))
def test_render_one_renders_each_layout_alone(runner, tmp_path, layout):
    out = tmp_path / f"{layout}.pptx"
    result = runner.invoke(cli, ["render-one", "--layout", layout, "--out", str(out)])
    assert result.exit_code == 0, result.output
    assert len(Presentation(str(out)).slides) == 1
    assert "Geometry check: clean." in result.output


def test_render_one_rejects_an_unknown_layout(runner, tmp_path):
    result = runner.invoke(cli, ["render-one", "--layout", "nope", "--out", str(tmp_path / "x")])
    assert result.exit_code != 0


def test_every_layout_has_a_sample_and_a_renderer():
    from deckpilot.renderer.deck import RENDERERS

    assert set(SAMPLES) <= set(RENDERERS)


# -- the working commands --------------------------------------------------


def test_layouts_lists_every_renderer(runner):
    from deckpilot.renderer.deck import RENDERERS

    result = runner.invoke(cli, ["layouts"])
    assert result.exit_code == 0, result.output
    for name in RENDERERS:
        assert name in result.output


def test_layouts_describes_each_one_from_its_own_docstring(runner):
    result = runner.invoke(cli, ["layouts"])
    assert "Contents: the sections, and where each one starts." in result.output
    assert "RAID log as a real PowerPoint table, grouped by type." in result.output


def test_check_passes_a_clean_deck(runner, tmp_path):
    out = tmp_path / "deck.pptx"
    assert runner.invoke(cli, ["demo", "--out", str(out)]).exit_code == 0
    result = runner.invoke(cli, ["check", str(out)])
    assert result.exit_code == 0
    assert "clean" in result.output


def test_check_refuses_a_file_that_is_not_there(runner, tmp_path):
    assert runner.invoke(cli, ["check", str(tmp_path / "nope.pptx")]).exit_code != 0


# -- speaker notes ---------------------------------------------------------


def test_content_slides_carry_speaker_notes(runner, tmp_path):
    out = tmp_path / "deck.pptx"
    assert runner.invoke(cli, ["demo", "--out", str(out)]).exit_code == 0

    prs = Presentation(str(out))
    noted = [s for s in prs.slides if s.has_notes_slide]
    assert len(noted) >= 8
    for slide in noted:
        assert slide.notes_slide.notes_text_frame.text.strip()


def test_dividers_carry_no_notes_page(runner, tmp_path):
    """Touching notes_slide creates one; a deck of empty notes looks like notes."""
    from deckpilot.data.generate import build_programme
    from deckpilot.renderer import deck as deck_module
    from deckpilot.specgen.fallback import build_deck_spec

    spec = build_deck_spec(build_programme())
    prs = deck_module.build(spec)
    for slide_spec, slide in zip(spec.slides, prs.slides, strict=True):
        if slide_spec.layout in ("section_divider", "agenda"):
            assert not slide.has_notes_slide, slide_spec.layout


def test_notes_quote_the_same_numbers_as_the_slide():
    """Notes written once and never updated are the usual failure; these are derived."""
    from deckpilot.data.generate import build_programme
    from deckpilot.specgen.fallback import build_deck_spec

    programme = build_programme()
    spec = build_deck_spec(programme)
    summary = next(s for s in spec.slides if s.layout == "exec_summary")
    assert summary.verdict in summary.notes

    scorecard = next(s for s in spec.slides if s.layout == "kpi_scorecard")
    behind = [r.name for r in scorecard.rows if r.attainment < r.expected]
    assert any(name in scorecard.notes for name in behind)


# -- the agenda ------------------------------------------------------------


def test_the_agenda_page_numbers_match_the_deck():
    from deckpilot.data.generate import build_programme
    from deckpilot.specgen.fallback import build_deck_spec

    spec = build_deck_spec(build_programme())
    agenda = next(s for s in spec.slides if s.layout == "agenda")
    dividers = {
        s.number: i for i, s in enumerate(spec.slides, start=1)
        if s.layout == "section_divider"
    }
    assert dividers
    for entry in agenda.entries:
        assert entry.page == dividers[entry.number], entry.title


def test_the_agenda_counts_the_whole_deck():
    from deckpilot.data.generate import build_programme
    from deckpilot.specgen.fallback import build_deck_spec

    spec = build_deck_spec(build_programme())
    agenda = next(s for s in spec.slides if s.layout == "agenda")
    assert f"{len(spec.slides)} slides" in agenda.subtitle


def test_the_agenda_survives_an_appendix_of_a_different_length():
    """Its numbers depend on how many slides the RAID log paginated into."""
    from deckpilot.data.generate import build_programme
    from deckpilot.data.models import Programme
    from deckpilot.specgen.fallback import build_deck_spec

    programme = build_programme()
    payload = programme.model_dump(mode="json")
    payload["raid"] = payload["raid"][:4]
    trimmed = Programme.model_validate(payload)

    for candidate in (programme, trimmed):
        spec = build_deck_spec(candidate)
        agenda = next(s for s in spec.slides if s.layout == "agenda")
        dividers = {
            s.number: i for i, s in enumerate(spec.slides, start=1)
            if s.layout == "section_divider"
        }
        for entry in agenda.entries:
            assert entry.page == dividers[entry.number]
