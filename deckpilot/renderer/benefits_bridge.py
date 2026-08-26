"""Cost bridge: where the baseline goes, and what moves it.

The exhibit a CFO looks for first, and the one that is easiest to draw
dishonestly. Two decisions matter.

The axis is truncated. A bridge from 42 to 33.5 drawn from zero turns every lever
into a sliver and the reader learns nothing; truncating makes the steps legible.
That is only acceptable if the slide says so, so it does.

And a step too small to see is still drawn, at a minimum height. Dropping it
would leave a gap in the bridge with no explanation, which is worse than a bar a
reader has to squint at - and the value label states the real number anyway.
"""

from __future__ import annotations

from pptx.enum.text import MSO_ANCHOR, PP_ALIGN
from pptx.presentation import Presentation as PresentationType
from pptx.slide import Slide

from deckpilot.renderer.base import (
    FitRequest,
    TextStyle,
    add_line,
    add_rect,
    add_slide,
    add_textbox,
    considerations_panel,
    fit_group,
    footer,
    title_block,
)
from deckpilot.specgen.schema import BenefitsBridgeSpec, BridgeStep
from deckpilot.theme import tokens as T

STEP_COLORS = {
    "anchor": T.PRIMARY,
    "decrease": T.STATUS_GREEN,  # takes cost out
    "increase": T.STATUS_RED,  # puts cost back
}


def axis_range(steps: list[BridgeStep]) -> tuple[float, float]:
    """Floor and ceiling for the value axis, truncated to keep steps legible."""
    levels = [v for step in steps for v in (step.from_value, step.to_value)]
    low, high = min(levels), max(levels)
    span = high - low or abs(high) or 1.0
    return low - span * T.BRIDGE_AXIS_PAD_BELOW, high + span * T.BRIDGE_AXIS_PAD_ABOVE


def _bar_span(step: BridgeStep, floor: float) -> tuple[float, float]:
    """The two values a column spans. An anchor stands on the axis floor."""
    if step.kind == "anchor":
        return floor, step.to_value
    return min(step.from_value, step.to_value), max(step.from_value, step.to_value)


def render(prs: PresentationType, spec: BenefitsBridgeSpec, page: int) -> Slide:
    slide = add_slide(prs)

    has_panel = bool(spec.considerations)
    total_w = T.body_width_with_panel() if has_panel else T.content_width()
    title_block(slide, spec.title, spec.subtitle, where=f"page {page}", width=total_w)

    n = len(spec.steps)
    bar_w = T.col_w(n, total=total_w, gutter=T.BRIDGE_GUTTER)

    chart_top = T.content_top() + T.BRIDGE_VALUE_BAND_H
    chart_bottom = T.content_bottom() - T.BRIDGE_LABEL_BAND_H
    chart_h = chart_bottom - chart_top

    floor, ceiling = axis_range(spec.steps)
    scale = chart_h / (ceiling - floor)

    def y_of(value: float) -> int:
        return int(chart_bottom - (value - floor) * scale)

    values, labels, captions = [], [], []
    for i, step in enumerate(spec.steps):
        x = T.col_x(i, n, total=total_w, gutter=T.BRIDGE_GUTTER)
        low, high = _bar_span(step, floor)
        top, bottom = y_of(high), y_of(low)
        height = max(T.BRIDGE_MIN_BAR_H, bottom - top)
        if height > bottom - top:
            # A widened sliver grows upward so it still meets its connector.
            top = bottom - height

        add_rect(slide, x, top, bar_w, height, fill=STEP_COLORS[step.kind], name=f"step{i}:bar")

        # The value sits above the column, in the band reserved for it.
        box = add_textbox(
            slide,
            x,
            top - T.BRIDGE_VALUE_BAND_H - T.BRIDGE_VALUE_GAP,
            bar_w,
            T.BRIDGE_VALUE_BAND_H,
            name=f"step{i}:value",
        )
        values.append(
            FitRequest(
                box.text_frame, [step.value], bar_w, T.BRIDGE_VALUE_BAND_H,
                TextStyle(
                    bold=True,
                    color=T.PRIMARY if step.kind == "anchor" else STEP_COLORS[step.kind],
                    align=PP_ALIGN.CENTER, anchor=MSO_ANCHOR.BOTTOM,
                    line_spacing=1.0, space_after_pt=0.0, wrap=False,
                ),
                f"bridge/{step.label}/value",
            )
        )

        caption_h = T.BRIDGE_CAPTION_H if step.caption else 0
        label = add_textbox(
            slide, x, chart_bottom, bar_w, T.BRIDGE_LABEL_BAND_H - caption_h,
            name=f"step{i}:label",
        )
        labels.append(
            FitRequest(
                label.text_frame, [step.label], bar_w,
                T.BRIDGE_LABEL_BAND_H - caption_h,
                TextStyle(
                    bold=step.kind == "anchor", color=T.GRAY_DARK, align=PP_ALIGN.CENTER,
                    anchor=MSO_ANCHOR.TOP, line_spacing=0.95, space_after_pt=0.0,
                ),
                f"bridge/{step.label}",
            )
        )
        if step.caption:
            caption = add_textbox(
                slide, x, chart_bottom + T.BRIDGE_LABEL_BAND_H - caption_h, bar_w, caption_h,
                name=f"step{i}:caption",
            )
            captions.append(
                FitRequest(
                    caption.text_frame, [step.caption], bar_w, caption_h,
                    TextStyle(
                        color=T.FOOTER_GRAY, align=PP_ALIGN.CENTER, anchor=MSO_ANCHOR.MIDDLE,
                        line_spacing=1.0, space_after_pt=0.0, wrap=False,
                    ),
                    f"bridge/{step.label}/caption",
                )
            )

        # Carry the level across to the next column so the bridge reads as one line.
        if i < n - 1:
            level = y_of(step.to_value)
            add_line(
                slide, x, level, T.col_x(i + 1, n, total=total_w, gutter=T.BRIDGE_GUTTER),
                level, color=T.tint(T.GRAY_DARK, 0.55), width_pt=T.BRIDGE_CONNECTOR_PT,
            )

    fit_group(values, T.FS_MICRO, T.FS_SUBTITLE)
    fit_group(labels, T.FS_MICRO - 1, T.FS_DENSE)
    if captions:
        fit_group(captions, T.FS_MICRO - 2, T.FS_MICRO - 1)

    if has_panel:
        considerations_panel(slide, spec.considerations, where=f"page {page}/considerations")

    footer(slide, page)
    return slide
