"""Workstream charter: the flagship layout.

N sub-stream columns, each with a numbered header, crossed by two labelled bands
that span the full width - key activities and outcomes. The bands are what make
it a charter rather than a set of lists: a reader can compare one row across all
columns without re-reading the headers.

This layout carries the densest text in the deck, so it gets the strictest
overflow treatment: every cell is fitted independently and any truncation is
reported by slide and by cell.
"""

from __future__ import annotations

from pptx.enum.shapes import MSO_SHAPE
from pptx.enum.text import MSO_ANCHOR, PP_ALIGN
from pptx.presentation import Presentation as PresentationType
from pptx.slide import Slide
from pptx.util import Emu

from deckpilot.renderer.base import (
    FitRequest,
    TextStyle,
    add_rect,
    add_slide,
    add_textbox,
    considerations_panel,
    fit_group,
    fit_text,
    footer,
    natural_height_pt,
    number_badge,
    set_vertical_text,
    title_block,
)
from deckpilot.specgen.schema import WorkstreamCharterSpec
from deckpilot.theme import tokens as T

# Whatever height a band does not need for text is spread across its bullet gaps
# by the group fit, rather than left at the bottom of the box.


def _list_icon(slide: Slide, x: int, y: int, color, key: str) -> None:
    """Three stacked rules - a list. Native shapes, so it stays editable."""
    bar_h = max(T.CHARTER_ICON_BAR_H, T.CHARTER_ICON_D // 8)
    step = T.CHARTER_ICON_D // 3
    for i in range(3):
        width = T.CHARTER_ICON_D if i < 2 else int(T.CHARTER_ICON_D * T.CHARTER_ICON_SHORT)
        add_rect(slide, x, y + i * step, width, bar_h, fill=color, name=f"{key}:list{i}")


def _target_icon(slide: Slide, x: int, y: int, color, key: str) -> None:
    """Two concentric rings and a centre - an outcome."""
    add_rect(
        slide, x, y, T.CHARTER_ICON_D, T.CHARTER_ICON_D, line=color, line_pt=T.HAIRLINE_PT,
        shape_type=MSO_SHAPE.OVAL, name=f"{key}:ring-outer",
    )
    inset = T.CHARTER_ICON_D // 3
    add_rect(
        slide, x + inset, y + inset, T.CHARTER_ICON_D - 2 * inset, T.CHARTER_ICON_D - 2 * inset,
        fill=color, shape_type=MSO_SHAPE.OVAL, name=f"{key}:ring-inner",
    )


def _band_label(slide: Slide, x: int, y: int, h: int, text: str, icon, key: str) -> None:
    """The vertical label strip: icon at the top, rotated caption beneath it."""
    add_rect(slide, x, y, T.CHARTER_LABEL_W, h, fill=T.tint(T.PRIMARY, 0.90), name=f"{key}:strip")
    icon_x = x + (T.CHARTER_LABEL_W - T.CHARTER_ICON_D) // 2
    icon(slide, icon_x, y + T.CHARTER_ICON_TOP_PAD, T.PRIMARY, key)

    caption_top = y + T.CHARTER_ICON_TOP_PAD + T.CHARTER_ICON_D + T.CHARTER_CAPTION_GAP
    caption_h = h - (caption_top - y) - T.TIGHT_GAP
    caption = add_textbox(
        slide, x, caption_top, T.CHARTER_LABEL_W, caption_h, name=f"{key}:caption"
    )
    set_vertical_text(caption.text_frame)
    fit_text(
        caption.text_frame,
        text,
        T.FS_MICRO - 1,
        T.FS_BODY,
        where=f"charter/{key} label",
        style=TextStyle(
            bold=True,
            color=T.PRIMARY,
            align=PP_ALIGN.RIGHT,  # with vert270 flow, this puts the text at the top
            anchor=MSO_ANCHOR.MIDDLE,
            line_spacing=1.0,
            space_after_pt=0.0,
            wrap=False,
        ),
        # The frame is rotated, so its usable extents are swapped.
        shape_w=caption_h,
        shape_h=T.CHARTER_LABEL_W,
    )
    caption.text_frame.margin_top = caption.text_frame.margin_bottom = Emu(0)


CELL_STYLE = TextStyle(
    color=T.GRAY_DARK,
    bullet="•",
    anchor=MSO_ANCHOR.TOP,
    line_spacing=T.LINE_SPACING,
    space_after_pt=T.SPACE_AFTER_PT,
)


def _cell(
    slide: Slide, x: int, y: int, w: int, h: int, items: list[str], where: str, key: str
) -> FitRequest:
    """Draw a cell and return its fit request; the band fits all of them together."""
    add_rect(slide, x, y, w, h, fill=T.WHITE, line=T.GRAY_LIGHT, name=f"{key}:cell")
    box = add_textbox(slide, x, y, w, h, name=f"{key}:text")
    return FitRequest(box.text_frame, items, w, h, CELL_STYLE, where)


def _band_heights(spec: WorkstreamCharterSpec, cell_w: int, available: int) -> tuple[int, int]:
    """Heights of the activities and outcomes bands, driven by their content."""
    size = T.FS_DENSE + 1
    tallest_a = max(
        natural_height_pt(c.activities, cell_w, size, CELL_STYLE) for c in spec.columns
    )
    tallest_o = max(
        natural_height_pt(c.outcomes, cell_w, size, CELL_STYLE) for c in spec.columns
    )

    usable = available - T.GUTTER
    share = tallest_a / (tallest_a + tallest_o)
    share = min(T.CHARTER_MAX_ACTIVITIES_SHARE, max(T.CHARTER_MIN_ACTIVITIES_SHARE, share))
    band_a = int(usable * share)
    return band_a, usable - band_a


def render(prs: PresentationType, spec: WorkstreamCharterSpec, page: int) -> Slide:
    slide = add_slide(prs)

    has_panel = bool(spec.considerations)
    total_w = T.body_width_with_panel() if has_panel else T.content_width()
    title_block(slide, spec.title, spec.subtitle, where=f"page {page}", width=total_w)

    # Columns start to the right of the vertical label strip.
    grid_x = T.content_left() + T.CHARTER_LABEL_W + T.GUTTER
    grid_w = total_w - T.CHARTER_LABEL_W - T.GUTTER
    n = len(spec.columns)
    cw = T.col_w(n, total=grid_w)

    top = T.content_top()
    bands_top = top + T.CHARTER_HEADER_H + T.GUTTER
    bands_h = T.content_height() - T.CHARTER_HEADER_H - T.GUTTER
    activities_h, outcomes_h = _band_heights(spec, cw, bands_h)
    outcomes_top = bands_top + activities_h + T.GUTTER

    # -- column headers ---------------------------------------------------
    header_requests: list[FitRequest] = []
    for i, column in enumerate(spec.columns):
        x = T.col_x(i, n, total=grid_w, x0=grid_x)
        add_rect(slide, x, top, cw, T.CHARTER_HEADER_H, fill=T.PRIMARY, name=f"col{i}:header")
        badge_y = top + (T.CHARTER_HEADER_H - T.BADGE_D) // 2
        number_badge(
            slide,
            x + T.CHARTER_BADGE_PAD,
            badge_y,
            column.number,
            fill=T.WHITE,
            text_color=T.PRIMARY,
            name=f"col{i}:badge",
        )
        name_x = x + T.CHARTER_BADGE_PAD + T.BADGE_D + T.LABEL_GAP
        name_w = x + cw - name_x - T.CHARTER_BADGE_PAD
        name = add_textbox(slide, name_x, top, name_w, T.CHARTER_HEADER_H, name=f"col{i}:name")
        header_requests.append(
            FitRequest(
                name.text_frame,
                [column.name],
                name_w,
                T.CHARTER_HEADER_H,
                TextStyle(
                    bold=True,
                    color=T.WHITE,
                    anchor=MSO_ANCHOR.MIDDLE,
                    line_spacing=0.95,
                    space_after_pt=0.0,
                ),
                f"page {page}/column {column.number} header",
            )
        )
    fit_group(header_requests, T.FS_MICRO, T.FS_COLUMN_HEADER)

    # -- bands ------------------------------------------------------------
    _band_label(
        slide, T.content_left(), bands_top, activities_h, spec.activities_label,
        _list_icon, "activities",
    )
    _band_label(
        slide, T.content_left(), outcomes_top, outcomes_h, spec.outcomes_label,
        _target_icon, "outcomes",
    )

    activity_cells: list[FitRequest] = []
    outcome_cells: list[FitRequest] = []
    for i, column in enumerate(spec.columns):
        x = T.col_x(i, n, total=grid_w, x0=grid_x)
        activity_cells.append(
            _cell(
                slide, x, bands_top, cw, activities_h, column.activities,
                f"page {page}/{column.name}/activities", f"col{i}act",
            )
        )
        outcome_cells.append(
            _cell(
                slide, x, outcomes_top, cw, outcomes_h, column.outcomes,
                f"page {page}/{column.name}/outcomes", f"col{i}out",
            )
        )
    # One type size per band, so peer cells never differ by a point.
    fit_group(activity_cells, T.FS_MIN_DENSE, T.FS_DENSE + 1, fill=True)
    fit_group(outcome_cells, T.FS_MIN_DENSE, T.FS_DENSE + 1, fill=True)

    if has_panel:
        considerations_panel(slide, spec.considerations, where=f"page {page}/considerations")

    footer(slide, page)
    return slide
