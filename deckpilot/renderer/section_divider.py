"""Section divider: a full-bleed break between parts of the deck.

No body content, no footer. The whole point is a pause, so the only things on it
are the section number, the section title and one geometric accent.
"""

from __future__ import annotations

from pptx.enum.shapes import MSO_SHAPE
from pptx.enum.text import MSO_ANCHOR, PP_ALIGN
from pptx.presentation import Presentation as PresentationType
from pptx.slide import Slide
from pptx.util import Emu

from deckpilot.renderer.base import TextStyle, add_rect, add_slide, add_textbox, fit_text
from deckpilot.specgen.schema import SectionDividerSpec
from deckpilot.theme import tokens as T

# The accent is a single oversized circle bled off the right edge. It is drawn in
# a lighter shade of the background rather than a new colour, so it reads as
# depth rather than decoration.
ACCENT_D = T.inches(7.2)
ACCENT_OFFSET_X = T.inches(1.6)  # how far past the right edge its centre sits
ACCENT_OFFSET_Y = T.inches(1.1)  # how far below the bottom edge its centre sits

NUMBER_W = T.inches(3.0)
# The box has to hold the numeral's full line box, or the fit shrinks the number
# below the oversized size this layout exists for. A numeral's cap height is well
# short of its em, so the gap down to the title reads a little larger than the
# measurement suggests; that is the price of not clipping the glyph.
NUMBER_H = T.inches(1.92)
TITLE_W = T.inches(7.4)
TITLE_H = T.inches(0.56)
TITLE_KICKER_GAP = T.inches(0.14)
KICKER_H = T.inches(0.38)
BLOCK_LEFT = T.inches(1.0)
BLOCK_TOP = T.inches(2.10)


def render(prs: PresentationType, spec: SectionDividerSpec) -> Slide:
    slide = add_slide(prs)

    add_rect(slide, 0, 0, T.SLIDE_W, T.SLIDE_H, fill=T.PRIMARY, name="bleed:background")

    accent_x = T.SLIDE_W - ACCENT_D + ACCENT_OFFSET_X
    accent_y = T.SLIDE_H - ACCENT_D + ACCENT_OFFSET_Y
    add_rect(
        slide,
        accent_x,
        accent_y,
        ACCENT_D,
        ACCENT_D,
        fill=T.tint(T.PRIMARY, 0.10),
        shape_type=MSO_SHAPE.OVAL,
        name="bleed:accent",
    )

    number = add_textbox(slide, BLOCK_LEFT, BLOCK_TOP, NUMBER_W, NUMBER_H, name="section:number")
    fit_text(
        number.text_frame,
        spec.number,
        T.FS_SECTION_NUMBER * 0.6,
        T.FS_SECTION_NUMBER,
        where=f"section {spec.number}/number",
        style=TextStyle(
            bold=True,
            color=T.tint(T.PRIMARY, 0.42),
            anchor=MSO_ANCHOR.MIDDLE,
            line_spacing=0.85,
            space_after_pt=0.0,
        ),
    )
    number.text_frame.margin_left = number.text_frame.margin_right = Emu(0)

    title_top = BLOCK_TOP + NUMBER_H
    title = add_textbox(slide, BLOCK_LEFT, title_top, TITLE_W, TITLE_H, name="section:title")
    fit_text(
        title.text_frame,
        spec.title,
        T.FS_SECTION_TITLE * 0.6,
        T.FS_SECTION_TITLE,
        where=f"section {spec.number}/title",
        style=TextStyle(
            bold=True,
            color=T.WHITE,
            anchor=MSO_ANCHOR.TOP,
            line_spacing=1.0,
            space_after_pt=0.0,
        ),
    )
    title.text_frame.margin_left = title.text_frame.margin_right = Emu(0)

    if spec.kicker:
        kicker = add_textbox(
            slide,
            BLOCK_LEFT,
            title_top + TITLE_H + TITLE_KICKER_GAP,
            TITLE_W,
            KICKER_H,
            name="section:kicker",
        )
        fit_text(
            kicker.text_frame,
            spec.kicker,
            T.FS_MICRO,
            T.FS_SUBTITLE,
            where=f"section {spec.number}/kicker",
            style=TextStyle(
                color=T.tint(T.PRIMARY, 0.62),
                align=PP_ALIGN.LEFT,
                anchor=MSO_ANCHOR.TOP,
                line_spacing=1.0,
                space_after_pt=0.0,
            ),
        )
        kicker.text_frame.margin_left = kicker.text_frame.margin_right = Emu(0)

    return slide
