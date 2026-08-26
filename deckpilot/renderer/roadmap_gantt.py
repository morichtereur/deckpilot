"""Roadmap: work package and sub-stream on the left, a month grid on the right.

The hard parts are all at the edges. A bar that starts before the window or ends
after it has to be cut at the grid edge, not drawn overhanging it. A one-week
task must stay visible without growing past the last month. A milestone on a bar
boundary has to sit on top of the bar rather than beside it. The arithmetic for
all of that lives in `timeline.py`; this module places shapes.
"""

from __future__ import annotations

from datetime import date

from pptx.enum.shapes import MSO_SHAPE
from pptx.enum.text import MSO_ANCHOR, PP_ALIGN
from pptx.presentation import Presentation as PresentationType
from pptx.slide import Slide
from pptx.util import Emu

from deckpilot.renderer import text_metrics as tm
from deckpilot.renderer.base import (
    FitRequest,
    TextStyle,
    add_line,
    add_rect,
    add_slide,
    add_textbox,
    considerations_panel,
    fit_group,
    fit_text,
    footer,
    title_block,
)
from deckpilot.renderer.timeline import MonthGrid, assign_lanes
from deckpilot.specgen.schema import GanttRow, RoadmapGanttSpec
from deckpilot.theme import tokens as T

BAR_LABEL_PT = T.FS_MICRO - 1  # bars are labelled at a fixed size or not at all


def _row_groups(rows: list[GanttRow]) -> list[tuple[str, int, int]]:
    """Consecutive rows sharing a work package, as (name, first index, count)."""
    groups: list[tuple[str, int, int]] = []
    for i, row in enumerate(rows):
        if groups and groups[-1][0] == row.work_package:
            name, first, count = groups[-1]
            groups[-1] = (name, first, count + 1)
        else:
            groups.append((row.work_package, i, 1))
    return groups


def _header(slide: Slide, x: int, y: int, w: int, grid: MonthGrid) -> None:
    add_rect(slide, x, y, w, T.GANTT_HEADER_H, fill=T.PRIMARY, name="header:band")

    style = TextStyle(
        bold=True,
        color=T.WHITE,
        anchor=MSO_ANCHOR.MIDDLE,
        line_spacing=1.0,
        space_after_pt=0.0,
        wrap=False,
    )
    requests: list[FitRequest] = []
    for label, cx, cw, align in (
        ("Work package", x, T.GANTT_WP_COL_W, PP_ALIGN.LEFT),
        ("Sub-stream", x + T.GANTT_WP_COL_W, T.GANTT_SS_COL_W, PP_ALIGN.LEFT),
    ):
        box = add_textbox(slide, cx, y, cw, T.GANTT_HEADER_H, name=f"header:{label}")
        requests.append(
            FitRequest(
                box.text_frame, [label], cw, T.GANTT_HEADER_H,
                TextStyle(**{**style.__dict__, "align": align}), f"roadmap/{label} header",
            )
        )

    month_w = int(grid.month_width)
    for i, label in enumerate(grid.labels()):
        box = add_textbox(
            slide, int(grid.month_x(i)), y, month_w, T.GANTT_HEADER_H, name=f"header:m{i}"
        )
        requests.append(
            FitRequest(
                box.text_frame, [label], month_w, T.GANTT_HEADER_H,
                TextStyle(**{**style.__dict__, "align": PP_ALIGN.CENTER}),
                f"roadmap/month {label}",
            )
        )
    fit_group(requests, T.FS_MICRO - 2, T.FS_MICRO)


def _labels(
    slide: Slide, rows: list[GanttRow], x: int, rows_top: int, row_h: int, tint_x: int, tint_w: int
) -> None:
    # Alternating tint runs from the sub-stream column across the grid, so the
    # work package blocks on the left can span several rows without colliding.
    for i in range(len(rows)):
        if i % 2:
            add_rect(
                slide, tint_x, rows_top + i * row_h, tint_w, row_h,
                fill=T.ROW_TINT, name=f"surface:row{i}tint",
            )

    wp_requests: list[FitRequest] = []
    for name, first, count in _row_groups(rows):
        y = rows_top + first * row_h
        h = count * row_h
        add_rect(
            slide, x, y, T.GANTT_WP_COL_W, h,
            fill=T.tint(T.PRIMARY, 0.90), name=f"wp{first}:block",
        )
        box = add_textbox(
            slide,
            x + T.GANTT_LABEL_PAD,
            y,
            T.GANTT_WP_COL_W - 2 * T.GANTT_LABEL_PAD,
            h,
            name=f"wp{first}:label",
        )
        wp_requests.append(
            FitRequest(
                box.text_frame, [name],
                T.GANTT_WP_COL_W - 2 * T.GANTT_LABEL_PAD, h,
                TextStyle(
                    bold=True, color=T.PRIMARY, anchor=MSO_ANCHOR.MIDDLE,
                    line_spacing=0.95, space_after_pt=0.0,
                ),
                f"roadmap/work package {name}",
            )
        )
    fit_group(wp_requests, T.FS_MICRO - 2, T.FS_DENSE)

    ss_requests: list[FitRequest] = []
    ss_x = x + T.GANTT_WP_COL_W + T.GANTT_LABEL_PAD
    ss_w = T.GANTT_SS_COL_W - 2 * T.GANTT_LABEL_PAD
    for i, row in enumerate(rows):
        box = add_textbox(
            slide, ss_x, rows_top + i * row_h, ss_w, row_h, name=f"row{i}:sublabel"
        )
        ss_requests.append(
            FitRequest(
                box.text_frame, [row.sub_stream], ss_w, row_h,
                TextStyle(
                    color=T.GRAY_DARK, anchor=MSO_ANCHOR.MIDDLE,
                    line_spacing=0.95, space_after_pt=0.0,
                ),
                f"roadmap/{row.sub_stream}",
            )
        )
    fit_group(ss_requests, T.FS_MICRO - 2, T.FS_DENSE)


def _grid_lines(slide: Slide, grid: MonthGrid, top: int, bottom: int) -> None:
    for i in range(1, grid.month_count):
        x = int(grid.month_x(i))
        add_line(slide, x, top, x, bottom, color=T.GRAY_LIGHT, width_pt=T.HAIRLINE_PT)
    add_line(slide, grid.right, top, grid.right, bottom, color=T.GRAY_LIGHT,
             width_pt=T.HAIRLINE_PT)


def _label_offset(
    row: GanttRow, grid: MonthGrid, x: int, w: int, label_w: int
) -> int | None:
    """Left inset at which a bar label clears every milestone diamond, or None.

    Milestones mark phase ends, which are also the next phase's start, so a naive
    rule drops most labels on the chart. Sliding the label to the right of any
    diamond that sits on it keeps the label wherever the bar is long enough.
    """
    half = T.GANTT_DIAMOND_D // 2
    marks = [grid.clamped(m.date) for m in row.milestones if grid.is_visible(m.date)]
    offset = T.GANTT_LABEL_PAD
    for _ in range(len(marks) + 1):
        start, end = x + offset, x + offset + label_w
        blocking = [m for m in marks if start - half <= m <= end + half]
        if not blocking:
            break
        offset = int(max(blocking) + half + T.GANTT_LABEL_PAD) - x
    else:
        return None
    if offset + label_w + T.GANTT_LABEL_PAD > w:
        return None
    return offset


def _bars(slide: Slide, rows: list[GanttRow], grid: MonthGrid, rows_top: int, row_h: int) -> None:
    for i, row in enumerate(rows):
        lane_of, lane_count = assign_lanes([(b.start, b.end) for b in row.bars])
        zone = int(
            row_h * (T.GANTT_BAR_H_RATIO if lane_count == 1 else T.GANTT_MULTILANE_H_RATIO)
        )
        bar_h = max(1, (zone - (lane_count - 1) * T.GANTT_LANE_GAP) // lane_count)
        zone_top = rows_top + i * row_h + (row_h - zone) // 2

        for j, bar in enumerate(row.bars):
            extent = grid.bar(bar.start, bar.end, min_width=T.GANTT_BAR_MIN_W)
            if extent is None:
                continue  # entirely outside the window
            x, w = extent
            y = zone_top + lane_of[j] * (bar_h + T.GANTT_LANE_GAP)
            shape = add_rect(
                slide, x, y, w, bar_h,
                fill=T.PHASE_COLORS[bar.status], rounded=True, name=f"row{i}:bar{j}",
            )
            if not bar.label or bar_h < T.GANTT_MIN_LABEL_BAR_H:
                continue  # a stacked lane is too shallow to carry type
            label_w = int(tm.text_width(bar.label, BAR_LABEL_PT) * T.EMU_PER_PT)
            offset = _label_offset(row, grid, x, w, label_w)
            if offset is None:
                continue  # too narrow to label; the row label already names the stream
            fit_text(
                shape.text_frame,
                bar.label,
                BAR_LABEL_PT,
                BAR_LABEL_PT,
                where=f"roadmap/{row.sub_stream}/{bar.label}",
                style=TextStyle(
                    color=T.on_color(T.PHASE_COLORS[bar.status]),
                    align=PP_ALIGN.LEFT,
                    anchor=MSO_ANCHOR.MIDDLE,
                    line_spacing=1.0,
                    space_after_pt=0.0,
                    wrap=False,
                ),
            )
            shape.text_frame.margin_left = Emu(offset)
            shape.text_frame.margin_right = Emu(0)
            shape.text_frame.margin_top = shape.text_frame.margin_bottom = Emu(0)


def _milestones(
    slide: Slide, rows: list[GanttRow], grid: MonthGrid, rows_top: int, row_h: int
) -> None:
    d = T.GANTT_DIAMOND_D
    for i, row in enumerate(rows):
        cy = rows_top + i * row_h + (row_h - d) // 2
        for j, milestone in enumerate(row.milestones):
            if not grid.is_visible(milestone.date):
                continue
            cx = int(grid.clamped(milestone.date)) - d // 2
            cx = min(max(cx, grid.x), grid.right - d)
            add_rect(
                slide, cx, cy, d, d,
                fill=T.PRIMARY if milestone.major else T.WHITE,
                line=None if milestone.major else T.PRIMARY,
                line_pt=T.HAIRLINE_PT,
                shape_type=MSO_SHAPE.DIAMOND,
                name=f"marker:row{i}ms{j}",
            )


def _today_line(slide: Slide, grid: MonthGrid, day: date, top: int, bottom: int) -> None:
    if not grid.is_visible(day):
        return
    x = int(grid.clamped(day))
    add_line(slide, x, top, x, bottom, color=T.SECONDARY, width_pt=T.TODAY_LINE_PT)


def _legend(slide: Slide, x: int, y: int, width: int) -> None:
    """Bar colours only mean something with a key, and the key also carries the
    two things that are not bars: the milestone marker and the reporting date."""
    entries: list[tuple[str, object]] = [
        (T.PHASE_LABELS[state], ("bar", T.PHASE_COLORS[state]))
        for state in ("complete", "in-progress", "planned", "at-risk")
    ]
    entries.append(("Milestone", ("diamond", T.PRIMARY)))
    entries.append(("Reporting date", ("line", T.SECONDARY)))

    cursor = x
    for i, (label, (kind, color)) in enumerate(entries):
        swatch_w = T.GANTT_LEGEND_SWATCH_W
        if kind == "bar":
            add_rect(
                slide, cursor, y + T.GANTT_LEGEND_H // 3, swatch_w, T.GANTT_LEGEND_H // 3,
                fill=color, rounded=True, name=f"legend{i}:swatch",
            )
        elif kind == "diamond":
            d = T.GANTT_LEGEND_H // 2
            add_rect(
                slide, cursor + (swatch_w - d) // 2, y + T.GANTT_LEGEND_H // 4, d, d,
                fill=color, shape_type=MSO_SHAPE.DIAMOND, name=f"legend{i}:swatch",
            )
        else:
            cx = cursor + swatch_w // 2
            add_line(
                slide, cx, y + T.GANTT_LEGEND_H // 5, cx, y + T.GANTT_LEGEND_H,
                color=color, width_pt=T.TODAY_LINE_PT,
            )
        text_x = cursor + swatch_w + T.GANTT_LABEL_PAD
        text_w = int(tm.text_width(label, T.FS_MICRO) * T.EMU_PER_PT) + T.GANTT_LEGEND_TEXT_PAD
        box = add_textbox(slide, text_x, y, text_w, T.GANTT_LEGEND_H, name=f"legend{i}:label")
        fit_text(
            box.text_frame,
            label,
            T.FS_MICRO - 1,
            T.FS_MICRO,
            where=f"roadmap/legend/{label}",
            style=TextStyle(
                color=T.FOOTER_GRAY,
                anchor=MSO_ANCHOR.MIDDLE,
                line_spacing=1.0,
                space_after_pt=0.0,
                wrap=False,
            ),
        )
        box.text_frame.margin_left = box.text_frame.margin_right = Emu(0)
        cursor = text_x + text_w + T.GANTT_LEGEND_ITEM_GAP
        if cursor > x + width:
            break


def render(prs: PresentationType, spec: RoadmapGanttSpec, page: int) -> Slide:
    slide = add_slide(prs)

    has_panel = bool(spec.considerations)
    total_w = T.body_width_with_panel() if has_panel else T.content_width()
    title_block(slide, spec.title, spec.subtitle, where=f"page {page}", width=total_w)

    x = T.content_left()
    grid_x = x + T.GANTT_WP_COL_W + T.GANTT_SS_COL_W
    grid_w = total_w - T.GANTT_WP_COL_W - T.GANTT_SS_COL_W
    grid = MonthGrid(spec.window_start, spec.window_end, grid_x, grid_w)

    header_top = T.content_top()
    rows_top = header_top + T.GANTT_HEADER_H
    legend_top = T.content_bottom() - T.GANTT_LEGEND_H
    rows_h = legend_top - T.GANTT_LEGEND_GAP - rows_top
    row_h = rows_h // len(spec.rows)
    rows_bottom = rows_top + row_h * len(spec.rows)

    _header(slide, x, header_top, total_w, grid)
    _labels(slide, spec.rows, x, rows_top, row_h, x + T.GANTT_WP_COL_W,
            T.GANTT_SS_COL_W + grid_w)
    _grid_lines(slide, grid, rows_top, rows_bottom)
    _bars(slide, spec.rows, grid, rows_top, row_h)
    _milestones(slide, spec.rows, grid, rows_top, row_h)
    if spec.today:
        _today_line(slide, grid, spec.today, rows_top, rows_bottom)
    _legend(slide, x, legend_top, total_w)

    if has_panel:
        considerations_panel(slide, spec.considerations, where=f"page {page}/considerations")

    footer(slide, page)
    return slide
