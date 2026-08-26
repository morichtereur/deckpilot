"""RAID log as a real PowerPoint table, grouped by type.

A native table rather than a grid of drawn rectangles: a reader who wants to add
a row, re-sort, or paste the log into a tracker should be able to. Severity is
carried by a filled cell rather than a separate chip shape, for the same reason -
one object, still coloured, still editable.
"""

from __future__ import annotations

import logging
from dataclasses import dataclass

from pptx.enum.text import MSO_ANCHOR, PP_ALIGN
from pptx.presentation import Presentation as PresentationType
from pptx.slide import Slide
from pptx.util import Emu, Pt

from deckpilot.renderer import text_metrics as tm
from deckpilot.renderer.base import (
    _drop_theme_style,
    add_slide,
    footer,
    title_block,
)
from deckpilot.specgen.schema import RaidRow, RaidTableSpec
from deckpilot.theme import tokens as T

log = logging.getLogger("deckpilot.renderer")

# Column widths as fractions of the content width. They sum to 1.
COLUMN_SHARE = {
    "severity": 0.045,
    "id": 0.052,
    "title": 0.300,
    "owner": 0.130,
    "due": 0.078,
    "mitigation": 0.395,
}
HEADINGS = ["", "ID", "Item", "Owner", "Due", "Mitigation and next step"]
GROUP_ORDER = ["risk", "issue", "dependency", "assumption"]


def column_widths(total: int) -> dict[str, int]:
    return {key: int(total * share) for key, share in COLUMN_SHARE.items()}


def _ordered(rows: list[RaidRow]) -> list[tuple[str, list[RaidRow]]]:
    groups = []
    for kind in GROUP_ORDER:
        members = [r for r in rows if r.kind == kind]
        if members:
            groups.append((kind, members))
    return groups


def _cell(cell, text: str, size: float, *, bold=False, color=None, align=PP_ALIGN.LEFT,
          fill=None) -> None:
    cell.margin_left = cell.margin_right = Emu(T.RAID_CELL_PAD)
    cell.margin_top = cell.margin_bottom = Emu(CELL_MARGIN)
    cell.vertical_anchor = MSO_ANCHOR.MIDDLE
    if fill is None:
        cell.fill.background()
    else:
        cell.fill.solid()
        cell.fill.fore_color.rgb = fill

    frame = cell.text_frame
    frame.word_wrap = True
    frame.clear()
    paragraph = frame.paragraphs[0]
    paragraph.alignment = align
    run = paragraph.add_run()
    run.text = text
    run.font.name = T.FONT
    run.font.size = Pt(size)
    run.font.bold = bold
    run.font.color.rgb = color if color is not None else T.GRAY_DARK


CELL_MARGIN = T.inches(0.03)
CANDIDATE_SIZES = [T.FS_DENSE, T.FS_DENSE - 0.5, T.FS_MICRO, T.FS_MICRO - 0.5, T.FS_MICRO - 1]
PREFERRED_SIZE = CANDIDATE_SIZES[0]
BALANCE_TOLERANCE = T.inches(0.02)
WIDTH_SAFETY = 0.96  # measure against slightly less width than the column really has
# A declared row height is a floor, not a ceiling: over-declare and the row is
# simply as tall as asked, under-declare and the renderer grows it and pushes the
# whole table down the slide. So both figures below err high on purpose, and the
# rows are never opened out to fill the leftover - an earlier version did that and
# the table walked into the footer, because filling to the last EMU leaves nothing
# for the gap between a calibrated figure and a real renderer.
ROW_OVERHEAD = T.inches(0.03)
# A table cell's line box is markedly taller than a text box's - measured at
# roughly 1.45x the font size against the 1.22x that text_metrics assumes for
# ordinary frames. Using the text-box figure here under-declares every row, and
# because a declared row height is a floor rather than a ceiling, the renderer
# quietly grows each one and walks the table off the bottom of the slide.
TABLE_LINE_HEIGHT = 1.48


def _row_lines(row: RaidRow, widths: dict[str, int], size: float) -> int:
    """How many lines the tallest cell in this row needs.

    Measured against slightly less than the real column width. Our Calibri
    metrics are close to a renderer's but not identical, and a table row height
    is a minimum: underestimate one cell by a single word and PowerPoint grows
    that row, which pushes every row below it down and the table off the slide.
    Erring wide costs a little whitespace; erring narrow costs the slide.
    """
    def usable(key: str) -> float:
        return (widths[key] - 2 * T.RAID_CELL_PAD) * WIDTH_SAFETY / T.EMU_PER_PT

    return max(
        tm.line_count(text, usable(key), size)
        for key, text in (
            ("title", row.title),
            ("mitigation", row.mitigation),
            ("owner", row.owner),
        )
    )


def _plan(rows: list[RaidRow], widths: dict[str, int], groups: int,
          available: int) -> tuple[float, list[int]]:
    """Choose a type size and a height for every row so the whole table fits.

    A row height set on a python-pptx table is a minimum, not a maximum: if a
    cell wraps to more lines than the height allows, PowerPoint grows the row and
    the table runs off the bottom of the slide. So the heights have to be
    computed from the wrapped line count rather than divided out evenly.
    """
    fixed = T.RAID_HEADER_H + groups * T.RAID_GROUP_H
    heights: list[int] = []
    for size in CANDIDATE_SIZES:
        line_h = size * TABLE_LINE_HEIGHT * T.EMU_PER_PT
        heights = [
            max(
                T.RAID_ROW_MIN_H,
                int(_row_lines(row, widths, size) * line_h) + 2 * CELL_MARGIN + ROW_OVERHEAD,
            )
            for row in rows
        ]
        if fixed + sum(heights) <= available:
            return size, heights
    log.warning(
        "raid table: %d rows do not fit %.2f in even at %.1fpt",
        len(rows), available / T.EMU_PER_INCH, CANDIDATE_SIZES[-1],
    )
    return CANDIDATE_SIZES[-1], heights


@dataclass(frozen=True)
class Page:
    """One slide's worth of a log too long for a single slide."""

    rows: list[RaidRow]
    continued_groups: list[str]


def _break_into_pages(
    rows: list[RaidRow],
    widths: dict[str, int],
    size: float,
    available: int,
    available_cap: int | None = None,
) -> list[Page]:
    """Greedily fill slides, repeating a group's header when it spans a break.

    `available` is the height budget to fill to; `available_cap` is the real
    slide height, so a single row taller than the budget still goes on a slide of
    its own rather than being deferred forever.
    """
    cap = available if available_cap is None else available_cap
    line_h = size * TABLE_LINE_HEIGHT * T.EMU_PER_PT

    def height(row: RaidRow) -> int:
        return max(
            T.RAID_ROW_MIN_H,
            int(_row_lines(row, widths, size) * line_h) + 2 * CELL_MARGIN + ROW_OVERHEAD,
        )

    pages: list[Page] = []
    current: list[RaidRow] = []
    continued: list[str] = []
    group: str | None = None
    used = T.RAID_HEADER_H

    for row in rows:
        changed = row.kind != group
        needed = (T.RAID_GROUP_H if changed else 0) + height(row)
        if current and used + needed > available:
            pages.append(Page(current, continued))
            # A group cut in half is re-announced on the next slide.
            continued = [] if changed else [row.kind]
            current = [row]
            group = row.kind
            used = T.RAID_HEADER_H + T.RAID_GROUP_H + height(row)
            continue
        if changed:
            used += T.RAID_GROUP_H
            group = row.kind
        used += height(row)
        current.append(row)

    if current:
        pages.append(Page(current, continued))
    for page in pages:
        if len(page.rows) == 1 and T.RAID_HEADER_H + T.RAID_GROUP_H + height(page.rows[0]) > cap:
            log.warning(
                "raid table: item %s does not fit a slide even alone", page.rows[0].id
            )
    return pages


def paginate(rows: list[RaidRow], total_width: int, available: int) -> list[Page]:
    """Split a log across as many slides as it needs, at the full dense size.

    Type size and page count trade against each other, and for an appendix the
    trade is one-sided: a slide costs nothing and an unreadable table costs the
    reader. So pagination measures at PREFERRED_SIZE and adds slides rather than
    shrinking type - the opposite of what a single summary slide does, where the
    slide count is fixed at one and the type has to give.

    Because PREFERRED_SIZE is also the largest size `_plan` will consider, every
    page produced here renders at exactly that size. That matters: a log whose
    second slide is set a point larger than its first looks like two different
    documents.
    """
    ordered = [row for _, members in _ordered(rows) for row in members]
    widths = column_widths(total_width)
    pages = _break_into_pages(ordered, widths, PREFERRED_SIZE, available)
    if len(pages) < 2:
        return pages

    # Filling greedily leaves the last slide with whatever is left over - sixteen
    # rows and then two. Shrinking the per-slide budget as far as it will go
    # without adding a slide spreads the same rows evenly across the same number
    # of slides, which is what someone laying this out by hand would do.
    target = len(pages)
    low, high = T.RAID_HEADER_H, available
    balanced = pages
    while high - low > BALANCE_TOLERANCE:
        middle = (low + high) // 2
        trial = _break_into_pages(
            ordered, widths, PREFERRED_SIZE, middle, available_cap=available
        )
        if len(trial) <= target:
            balanced, high = trial, middle
        else:
            low = middle + 1
    return balanced


def render(prs: PresentationType, spec: RaidTableSpec, page: int) -> Slide:
    slide = add_slide(prs)
    title_block(slide, spec.title, spec.subtitle, where=f"page {page}")

    groups = _ordered(spec.rows)
    total_rows = 1 + sum(1 + len(members) for _, members in groups)

    left, top = T.content_left(), T.content_top()
    width, height = T.content_width(), T.content_height()

    graphic = slide.shapes.add_table(total_rows, len(HEADINGS), Emu(left), Emu(top),
                                     Emu(width), Emu(height))
    graphic.name = "raid:table"
    table = graphic.table
    # python-pptx applies a banded theme style; every fill here is set explicitly.
    table.first_row = False
    table.horz_banding = False
    _drop_theme_style(graphic)

    widths = column_widths(width)
    for i, key in enumerate(COLUMN_SHARE):
        table.columns[i].width = Emu(widths[key])

    size, row_heights = _plan(spec.rows, widths, len(groups), height)
    # The frame is sized to what the rows actually need, so the geometry check
    # sees the same table the renderer will draw.
    graphic.height = Emu(T.RAID_HEADER_H + len(groups) * T.RAID_GROUP_H + sum(row_heights))

    table.rows[0].height = Emu(T.RAID_HEADER_H)
    for i, heading in enumerate(HEADINGS):
        _cell(table.cell(0, i), heading, T.FS_MICRO, bold=True, color=T.WHITE, fill=T.PRIMARY)

    index = 1
    for kind, members in groups:
        table.rows[index].height = Emu(T.RAID_GROUP_H)
        label = spec.group_labels.get(kind, kind.title())
        # A count on a continued group would be the count on this slide, not the
        # count in the log, which is worse than no count at all.
        label = (
            f"{label} (continued)"
            if kind in spec.continued_groups
            else f"{label} ({len(members)})"
        )
        _cell(table.cell(index, 0), "", T.FS_MICRO, fill=T.tint(T.PRIMARY, 0.86))
        _cell(table.cell(index, 1), label, T.FS_MICRO,
              bold=True, color=T.PRIMARY, fill=T.tint(T.PRIMARY, 0.86))
        table.cell(index, 1).merge(table.cell(index, len(HEADINGS) - 1))
        index += 1

        for row in members:
            table.rows[index].height = Emu(row_heights[spec.rows.index(row)])
            tint = T.WHITE if index % 2 else T.ROW_TINT
            severity_fill = T.SEVERITY_COLORS[row.severity]
            _cell(table.cell(index, 0), row.severity, size, bold=True,
                  color=T.on_color(severity_fill), align=PP_ALIGN.CENTER, fill=severity_fill)
            for column, text, align in (
                (1, row.id, PP_ALIGN.LEFT),
                (2, row.title, PP_ALIGN.LEFT),
                (3, row.owner, PP_ALIGN.LEFT),
                (4, row.due, PP_ALIGN.LEFT),
                (5, row.mitigation, PP_ALIGN.LEFT),
            ):
                _cell(table.cell(index, column), text, size, align=align, fill=tint)
            index += 1

    footer(slide, page)
    return slide
