"""The one slide that gets read: the verdict, the messages behind it, the asks.

Laid out so it can be read top to bottom in that order - a rated verdict band, a
row of key messages each carrying its own rating, and the decisions the meeting
is actually being asked for.
"""

from __future__ import annotations

from pptx.enum.text import MSO_ANCHOR, PP_ALIGN
from pptx.presentation import Presentation as PresentationType
from pptx.slide import Slide
from pptx.util import Emu

from deckpilot.renderer.base import (
    FitRequest,
    TextStyle,
    add_rect,
    add_slide,
    add_textbox,
    fit_group,
    fit_text,
    footer,
    natural_height_pt,
    title_block,
)
from deckpilot.specgen.schema import ExecSummarySpec
from deckpilot.theme import tokens as T

# The message cards are sized to what they hold; the decisions band takes
# whatever is left. A key message is two or three lines, and a card five times
# that tall reads as a card someone forgot to fill in.
MESSAGE_MIN_H = T.inches(1.05)
MESSAGE_MAX_SHARE = 0.62  # of the height below the verdict band
DECISIONS_MIN_H = T.inches(0.80)
DETAIL_STYLE = TextStyle(
    color=T.GRAY_DARK, anchor=MSO_ANCHOR.TOP, line_spacing=T.LINE_SPACING, space_after_pt=0.0
)
DECISIONS_STYLE = TextStyle(
    color=T.GRAY_DARK,
    bullet="•",
    anchor=MSO_ANCHOR.TOP,
    line_spacing=T.LINE_SPACING,
    space_after_pt=T.SPACE_AFTER_PT + 1,
)
DECISIONS_LABEL_H = T.inches(0.22)


def _verdict_band(slide: Slide, spec: ExecSummarySpec, x: int, y: int, w: int, h: int) -> None:
    colour = T.RAG_COLORS[spec.overall_rag]
    add_rect(slide, x, y, w, h, fill=T.PANEL_BG, name="verdict:box")
    add_rect(slide, x, y, T.EXEC_RAG_W, h, fill=colour, name="verdict:rag")

    rag = add_textbox(slide, x, y, T.EXEC_RAG_W, h, name="verdict:ragtext")
    fit_text(
        rag.text_frame,
        spec.overall_rag.upper(),
        T.FS_BODY,
        T.FS_SUBTITLE + 2,
        where="exec summary/overall rating",
        style=TextStyle(
            bold=True, color=T.on_color(colour), align=PP_ALIGN.CENTER,
            anchor=MSO_ANCHOR.MIDDLE, line_spacing=1.0, space_after_pt=0.0, wrap=False,
        ),
    )
    rag.text_frame.margin_left = rag.text_frame.margin_right = Emu(0)

    text_x = x + T.EXEC_RAG_W + T.EXEC_PAD
    verdict = add_textbox(
        slide, text_x, y, x + w - text_x - T.EXEC_PAD, h, name="verdict:text"
    )
    fit_text(
        verdict.text_frame,
        spec.verdict,
        T.FS_BODY,
        T.FS_SUBTITLE + 2,
        where="exec summary/verdict",
        style=TextStyle(
            bold=True, color=T.PRIMARY, anchor=MSO_ANCHOR.MIDDLE,
            line_spacing=1.0, space_after_pt=0.0,
        ),
    )


def render(prs: PresentationType, spec: ExecSummarySpec, page: int) -> Slide:
    slide = add_slide(prs)
    title_block(slide, spec.title, spec.subtitle, where=f"page {page}")

    x, w = T.content_left(), T.content_width()
    top, height = T.content_top(), T.content_height()

    _verdict_band(slide, spec, x, top, w, T.EXEC_VERDICT_H)

    rest_top = top + T.EXEC_VERDICT_H + T.EXEC_GAP
    rest_h = height - T.EXEC_VERDICT_H - T.EXEC_GAP

    n = len(spec.messages)
    cw = T.col_w(n, total=w)
    inner_message_w = cw - 2 * T.EXEC_PAD
    tallest = max(
        natural_height_pt([m.detail], inner_message_w, T.FS_BODY, DETAIL_STYLE)
        for m in spec.messages
    )
    wanted = (
        int(tallest * T.EMU_PER_PT) + T.EXEC_MESSAGE_HEAD_H + 2 * T.EXEC_PAD
    )
    if spec.decisions:
        messages_h = max(MESSAGE_MIN_H, min(wanted, int(rest_h * MESSAGE_MAX_SHARE)))
        # The decisions band hugs its bullets too. Every block on this slide is
        # sized to its content and the page simply ends where the content does;
        # a half-empty grey box is worse than white space.
        needed = natural_height_pt(
            list(spec.decisions), w - 2 * T.EXEC_PAD, T.FS_BODY, DECISIONS_STYLE
        )
        decisions_h = int(needed * T.EMU_PER_PT) + DECISIONS_LABEL_H + T.EXEC_PAD
        decisions_h = max(DECISIONS_MIN_H, min(decisions_h, rest_h - messages_h - T.EXEC_GAP))
    else:
        messages_h = rest_h
        decisions_h = 0
    headings, details = [], []
    for i, message in enumerate(spec.messages):
        mx = T.col_x(i, n, total=w)
        colour = T.RAG_COLORS[message.rag]
        add_rect(slide, mx, rest_top, cw, messages_h, fill=T.WHITE, line=T.GRAY_LIGHT,
                 name=f"msg{i}:box")
        add_rect(slide, mx, rest_top, cw, T.EXEC_MESSAGE_HEAD_H, fill=colour,
                 name=f"msg{i}:head")

        inner_x = mx + T.EXEC_PAD
        inner_w = inner_message_w
        heading = add_textbox(
            slide, inner_x, rest_top, inner_w, T.EXEC_MESSAGE_HEAD_H, name=f"msg{i}:heading"
        )
        headings.append(
            FitRequest(
                heading.text_frame, [message.heading], inner_w, T.EXEC_MESSAGE_HEAD_H,
                TextStyle(bold=True, color=T.on_color(colour), anchor=MSO_ANCHOR.MIDDLE,
                          line_spacing=0.95, space_after_pt=0.0),
                f"exec summary/{message.heading}",
            )
        )
        detail_y = rest_top + T.EXEC_MESSAGE_HEAD_H + T.EXEC_PAD
        detail_h = rest_top + messages_h - detail_y - T.EXEC_PAD
        detail = add_textbox(slide, inner_x, detail_y, inner_w, detail_h, name=f"msg{i}:detail")
        details.append(
            FitRequest(
                detail.text_frame, [message.detail], inner_w, detail_h,
                DETAIL_STYLE,
                f"exec summary/{message.heading}/detail",
            )
        )

    fit_group(headings, T.FS_MICRO, T.FS_COLUMN_HEADER)
    fit_group(details, T.FS_MIN_BODY, T.FS_BODY)

    if spec.decisions:
        dy = rest_top + messages_h + T.EXEC_GAP
        add_rect(slide, x, dy, w, decisions_h, fill=T.PANEL_BG, name="decisions:box")
        label_h = DECISIONS_LABEL_H
        label = add_textbox(
            slide, x + T.EXEC_PAD, dy + T.EXEC_PAD // 2, w - 2 * T.EXEC_PAD, label_h,
            name="decisions:label",
        )
        fit_text(
            label.text_frame,
            spec.decisions_label.upper(),
            T.FS_MICRO - 1,
            T.FS_DENSE,
            where="exec summary/decisions label",
            style=TextStyle(
                bold=True, color=T.SECONDARY, anchor=MSO_ANCHOR.MIDDLE,
                line_spacing=1.0, space_after_pt=0.0, wrap=False,
            ),
        )
        body_y = dy + T.EXEC_PAD // 2 + label_h
        body = add_textbox(
            slide, x + T.EXEC_PAD, body_y, w - 2 * T.EXEC_PAD,
            dy + decisions_h - body_y - T.EXEC_PAD // 2, name="decisions:body",
        )
        fit_text(
            body.text_frame,
            list(spec.decisions),
            T.FS_MIN_DENSE,
            T.FS_BODY,
            where="exec summary/decisions",
            style=DECISIONS_STYLE,
        )

    footer(slide, page)
    return slide
