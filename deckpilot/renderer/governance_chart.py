"""Governance: who decides, who does the work, and who has to be consulted.

Three tiers - steering committee, programme management, then one box per work
package split into the team that does the work and the teams that have to be
brought along. The connectors are native line shapes rather than a drawn image,
so the chart survives being edited when the governance changes, which it will.
"""

from __future__ import annotations

from pptx.enum.text import MSO_ANCHOR, PP_ALIGN
from pptx.presentation import Presentation as PresentationType
from pptx.slide import Slide

from deckpilot.renderer.base import (
    FitRequest,
    TextStyle,
    add_line,
    add_rect,
    add_slide,
    add_textbox,
    considerations_panel,
    fit_group,
    fit_text,
    footer,
    number_badge,
    title_block,
)
from deckpilot.specgen.schema import GovernanceBox, GovernanceChartSpec, GovernanceUnit
from deckpilot.theme import tokens as T

MEMBER_STYLE = TextStyle(
    color=T.GRAY_DARK,
    anchor=MSO_ANCHOR.TOP,
    line_spacing=1.0,
    space_after_pt=1.5,
)
TEAM_STYLE = TextStyle(
    color=T.GRAY_DARK,
    bullet="–",  # an en dash: a team list is not a list of actions
    bullet_indent=T.BULLET_INDENT,
    anchor=MSO_ANCHOR.TOP,
    line_spacing=1.0,
    space_after_pt=2.0,
)


def _panelled_box(slide: Slide, x: int, y: int, w: int, h: int, title: str,
                  caption: str | None, key: str) -> None:
    """A bordered white box with a filled header strip, as used by every tier."""
    add_rect(slide, x, y, w, h, fill=T.WHITE, line=T.GRAY_LIGHT, name=f"{key}:box")
    add_rect(slide, x, y, w, T.GOV_BOX_HEAD_H, fill=T.PRIMARY, name=f"{key}:head")

    caption_w = int(w * T.GOV_CAPTION_SHARE) if caption else 0
    title_w = w - caption_w - 2 * T.GOV_PAD
    head = add_textbox(
        slide, x + T.GOV_PAD, y, title_w, T.GOV_BOX_HEAD_H, name=f"{key}:title"
    )
    fit_text(
        head.text_frame,
        title,
        T.FS_MICRO - 1,
        T.FS_BODY,
        where=f"governance/{title}",
        style=TextStyle(
            bold=True, color=T.WHITE, anchor=MSO_ANCHOR.MIDDLE,
            line_spacing=1.0, space_after_pt=0.0, wrap=False,
        ),
    )
    if caption:
        cap = add_textbox(
            slide, x + T.GOV_PAD + title_w, y, caption_w - T.GOV_PAD,
            T.GOV_BOX_HEAD_H, name=f"{key}:caption",
        )
        fit_text(
            cap.text_frame,
            caption,
            T.FS_MICRO - 2,
            T.FS_MICRO,
            where=f"governance/{title} cadence",
            style=TextStyle(
                color=T.tint(T.PRIMARY, 0.62), align=PP_ALIGN.RIGHT,
                anchor=MSO_ANCHOR.MIDDLE, line_spacing=1.0, space_after_pt=0.0, wrap=False,
            ),
        )


def _tier(slide: Slide, box: GovernanceBox, x: int, y: int, w: int, h: int,
          key: str) -> list[FitRequest]:
    """Steering committee or programme management: a box of named people.

    Six names stacked in one column would make the box twice as tall as it needs
    to be, so anything over three splits into two columns.
    """
    _panelled_box(slide, x, y, w, h, box.title, box.caption, key)

    columns = 2 if len(box.members) > 3 else 1
    per_column = -(-len(box.members) // columns)
    body_y = y + T.GOV_BOX_HEAD_H + T.GOV_PAD
    body_h = h - (body_y - y) - T.GOV_PAD
    inner_w = w - 2 * T.GOV_PAD

    requests: list[FitRequest] = []
    for c in range(columns):
        chunk = box.members[c * per_column : (c + 1) * per_column]
        if not chunk:
            continue
        cw = T.col_w(columns, total=inner_w)
        cx = T.col_x(c, columns, total=inner_w, x0=x + T.GOV_PAD)
        cell = add_textbox(slide, cx, body_y, cw, body_h, name=f"{key}:members{c}")
        requests.append(
            FitRequest(cell.text_frame, chunk, cw, body_h, MEMBER_STYLE, f"governance/{box.title}")
        )
    return requests


def _unit(slide: Slide, unit: GovernanceUnit, x: int, y: int, w: int, h: int,
          core_label: str, contributing_label: str, key: str) -> tuple[list, list, list]:
    """One work package: who does it, and who has to be brought along."""
    add_rect(slide, x, y, w, h, fill=T.WHITE, line=T.GRAY_LIGHT, name=f"{key}:box")
    head_h = T.GOV_BOX_HEAD_H + T.GOV_HEAD_EXTRA
    add_rect(slide, x, y, w, head_h, fill=T.PRIMARY, name=f"{key}:head")

    badge_y = y + (head_h - T.BADGE_D) // 2
    number_badge(
        slide, x + T.GOV_PAD, badge_y, unit.number,
        fill=T.WHITE, text_color=T.PRIMARY, name=f"{key}:badge",
    )
    name_x = x + T.GOV_PAD + T.BADGE_D + T.TIGHT_GAP
    name_w = x + w - name_x - T.GOV_PAD
    name = add_textbox(slide, name_x, y, name_w, head_h, name=f"{key}:name")
    header_request = FitRequest(
        name.text_frame, [unit.name], name_w, head_h,
        TextStyle(
            bold=True, color=T.WHITE, anchor=MSO_ANCHOR.MIDDLE,
            line_spacing=0.95, space_after_pt=0.0,
        ),
        f"governance/{unit.name}",
    )

    inner_x = x + T.GOV_PAD
    inner_w = w - 2 * T.GOV_PAD
    body_top = y + head_h + T.GOV_PAD
    body_h = h - (body_top - y) - T.GOV_PAD

    # The two sections split the body in proportion to how many rows each holds.
    weight = len(unit.core_team) / (len(unit.core_team) + len(unit.contributing_teams))
    section_h = body_h - 2 * T.GOV_LABEL_H - T.GOV_SECTION_GAP
    core_h = int(section_h * weight)
    contributing_h = section_h - core_h

    sections = []
    cursor = body_top
    for label, items, height in (
        (core_label, unit.core_team, core_h),
        (contributing_label, unit.contributing_teams, contributing_h),
    ):
        tag = add_textbox(slide, inner_x, cursor, inner_w, T.GOV_LABEL_H, name=f"{key}:{label}")
        fit_text(
            tag.text_frame,
            label.upper(),
            T.FS_MICRO - 2,
            T.FS_MICRO - 1,
            where=f"governance/{unit.name}/{label}",
            style=TextStyle(
                bold=True, color=T.SECONDARY, anchor=MSO_ANCHOR.MIDDLE,
                line_spacing=1.0, space_after_pt=0.0, wrap=False,
            ),
        )
        cursor += T.GOV_LABEL_H
        cell = add_textbox(slide, inner_x, cursor, inner_w, height, name=f"{key}:{label}list")
        sections.append(
            FitRequest(
                cell.text_frame, items, inner_w, height, TEAM_STYLE,
                f"governance/{unit.name}/{label}",
            )
        )
        cursor += height + T.GOV_SECTION_GAP

    return [header_request], [sections[0]], [sections[1]]


def _connect(slide: Slide, steering_bottom: int, management: tuple[int, int, int],
             unit_tops: list[tuple[int, int]], bus_y: int) -> None:
    """Native lines, so the chart stays editable when the governance changes."""
    mgmt_x, mgmt_top, mgmt_bottom = management
    add_line(slide, mgmt_x, steering_bottom, mgmt_x, mgmt_top,
             color=T.SECONDARY, width_pt=T.GOV_CONNECTOR_PT)
    add_line(slide, mgmt_x, mgmt_bottom, mgmt_x, bus_y,
             color=T.SECONDARY, width_pt=T.GOV_CONNECTOR_PT)
    add_line(slide, unit_tops[0][0], bus_y, unit_tops[-1][0], bus_y,
             color=T.SECONDARY, width_pt=T.GOV_CONNECTOR_PT)
    for cx, top in unit_tops:
        add_line(slide, cx, bus_y, cx, top, color=T.SECONDARY, width_pt=T.GOV_CONNECTOR_PT)


def render(prs: PresentationType, spec: GovernanceChartSpec, page: int) -> Slide:
    slide = add_slide(prs)

    has_panel = bool(spec.considerations)
    total_w = T.body_width_with_panel() if has_panel else T.content_width()
    title_block(slide, spec.title, spec.subtitle, where=f"page {page}", width=total_w)

    left = T.content_left()
    centre = left + total_w // 2
    top = T.content_top()

    steering_x = centre - T.GOV_STEERING_W // 2
    steering_requests = _tier(
        slide, spec.steering, steering_x, top, T.GOV_STEERING_W, T.GOV_STEERING_H, "steerco"
    )

    mgmt_top = top + T.GOV_STEERING_H + T.GOV_DROP
    mgmt_x = centre - T.GOV_MANAGEMENT_W // 2
    mgmt_requests = _tier(
        slide, spec.programme_management, mgmt_x, mgmt_top,
        T.GOV_MANAGEMENT_W, T.GOV_MANAGEMENT_H, "pmo",
    )
    # One size across both tiers: they are peers and must read as peers.
    fit_group(steering_requests + mgmt_requests, T.FS_MICRO - 2, T.FS_MICRO)

    mgmt_bottom = mgmt_top + T.GOV_MANAGEMENT_H
    bus_y = mgmt_bottom + T.GOV_BUS_DROP
    units_top = bus_y + T.GOV_BUS_DROP
    units_h = T.content_bottom() - units_top

    n = len(spec.units)
    uw = T.col_w(n, total=total_w)
    headers, cores, contributors = [], [], []
    unit_tops = []
    for i, unit in enumerate(spec.units):
        ux = T.col_x(i, n, total=total_w, x0=left)
        h, c, x = _unit(
            slide, unit, ux, units_top, uw, units_h,
            spec.core_label, spec.contributing_label, f"unit{i}",
        )
        headers += h
        cores += c
        contributors += x
        unit_tops.append((ux + uw // 2, units_top))

    fit_group(headers, T.FS_MICRO - 1, T.FS_COLUMN_HEADER)
    # Core and contributing lists are fitted together so the two halves of a box
    # do not end up a point apart from each other.
    fit_group(cores + contributors, T.FS_MICRO - 2, T.FS_DENSE, fill=True)

    _connect(slide, top + T.GOV_STEERING_H, (centre, mgmt_top, mgmt_bottom), unit_tops, bus_y)

    if has_panel:
        considerations_panel(
            slide, spec.considerations, heading="Comments",
            where=f"page {page}/comments",
        )

    footer(slide, page)
    return slide
