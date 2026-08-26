"""Programme budget: what was approved, what has been spent, where it lands.

The variance column is a diverging bar about a zero line rather than a number in
a coloured cell. A reader scanning a budget wants to know which line is the
problem before they read any figure, and a bar answers that at a glance while a
red number still has to be read.

Beneath it sit the gauges - the contingency drawn against the programme elapsed.
Total forecast equal to total budget looks like control; contingency running down
faster than the calendar is what that control is costing.
"""

from __future__ import annotations

from pptx.enum.text import MSO_ANCHOR, PP_ALIGN
from pptx.presentation import Presentation as PresentationType
from pptx.slide import Slide

from deckpilot.renderer.base import (
    FitRequest,
    TextStyle,
    add_line,
    add_rect,
    add_slide,
    add_textbox,
    considerations_panel,
    fit_group,
    footer,
    title_block,
)
from deckpilot.specgen.schema import CostRow, FinancialSummarySpec, Gauge
from deckpilot.theme import tokens as T

COLUMN_SHARE = {
    "category": 0.255,
    "budget": 0.090,
    "actual": 0.110,
    "forecast": 0.095,
    "variance": 0.090,
    "bar": 0.360,
}
HEADINGS = {
    "category": "Cost category",
    "budget": "Budget",
    "actual": "Actual to date",
    "forecast": "Forecast at completion",
    "variance": "Variance",
    "bar": "Forecast against budget",
}
NUMERIC = ("budget", "actual", "forecast", "variance")


def variance_colour(value: float):
    """Overruns read red, underspend green, on-plan neutral."""
    if value > 0:
        return T.STATUS_RED
    if value < 0:
        return T.STATUS_GREEN
    return T.STATUS_NEUTRAL


def _columns(x: int, width: int) -> dict[str, tuple[int, int]]:
    out, cursor = {}, x
    for key, share in COLUMN_SHARE.items():
        w = int(width * share)
        out[key] = (cursor, w)
        cursor += w
    return out


def _variance_bar(slide: Slide, row: CostRow, x: int, y: int, w: int, scale: float,
                  key: str) -> None:
    """A bar either side of a zero line, so the problem line is visible before it is read."""
    centre = x + w // 2
    add_line(
        slide, centre, y, centre, y + T.FIN_BAR_H,
        color=T.tint(T.GRAY_DARK, 0.55), width_pt=T.FIN_ZERO_LINE_PT,
    )
    if row.variance_value == 0:
        return
    length = max(T.FIN_MIN_BAR_W, int(abs(row.variance_value) * scale))
    left = centre if row.variance_value > 0 else centre - length
    add_rect(
        slide, left, y, length, T.FIN_BAR_H,
        fill=variance_colour(row.variance_value), name=f"marker:{key}var",
    )


def _gauges(slide: Slide, gauges: list[Gauge], x: int, y: int, w: int) -> list[FitRequest]:
    """Two short bars on a shared scale: what is consumed against what has elapsed."""
    requests: list[FitRequest] = []
    track_x = x + T.FIN_GAUGE_LABEL_W
    track_w = w - T.FIN_GAUGE_LABEL_W - T.FIN_GAUGE_VALUE_W - T.FIN_PAD

    for i, gauge in enumerate(gauges):
        top = y + i * (T.FIN_GAUGE_H + T.FIN_GAUGE_GAP)
        label = add_textbox(
            slide, x, top, T.FIN_GAUGE_LABEL_W, T.FIN_GAUGE_H, name=f"gauge{i}:label"
        )
        requests.append(
            FitRequest(
                label.text_frame, [gauge.label], T.FIN_GAUGE_LABEL_W, T.FIN_GAUGE_H,
                TextStyle(
                    color=T.GRAY_DARK, anchor=MSO_ANCHOR.MIDDLE,
                    line_spacing=1.0, space_after_pt=0.0, wrap=False,
                ),
                f"financials/{gauge.label}",
            )
        )
        add_rect(slide, track_x, top, track_w, T.FIN_GAUGE_H, fill=T.GRAY_LIGHT,
                 name=f"gauge{i}:track")
        filled = int(track_w * gauge.fraction)
        if filled:
            add_rect(
                slide, track_x, top, filled, T.FIN_GAUGE_H,
                fill=T.SECONDARY if i else T.PRIMARY, name=f"marker:gauge{i}fill",
            )
        value = add_textbox(
            slide, track_x + track_w + T.FIN_PAD, top, T.FIN_GAUGE_VALUE_W, T.FIN_GAUGE_H,
            name=f"gauge{i}:value",
        )
        requests.append(
            FitRequest(
                value.text_frame, [gauge.value], T.FIN_GAUGE_VALUE_W, T.FIN_GAUGE_H,
                TextStyle(
                    bold=True, color=T.PRIMARY, align=PP_ALIGN.RIGHT,
                    anchor=MSO_ANCHOR.MIDDLE, line_spacing=1.0, space_after_pt=0.0, wrap=False,
                ),
                f"financials/{gauge.label}/value",
            )
        )
    return requests


def _row(slide: Slide, row: CostRow, columns, y: int, row_h: int, scale: float, key: str,
         requests: dict[str, list[FitRequest]]) -> None:
    for name, (cx, cw) in columns.items():
        if name == "bar":
            _variance_bar(
                slide, row, cx, y + (row_h - T.FIN_BAR_H) // 2, cw, scale, key
            )
            continue
        text = getattr(row, name)
        colour = T.GRAY_DARK
        if name == "variance" and row.variance_value:
            colour = variance_colour(row.variance_value)
        elif row.emphasis:
            colour = T.PRIMARY
        box = add_textbox(
            slide, cx + T.FIN_PAD, y, cw - 2 * T.FIN_PAD, row_h, name=f"{key}:{name}"
        )
        requests["category" if name == "category" else "figure"].append(
            FitRequest(
                box.text_frame, [text], cw - 2 * T.FIN_PAD, row_h,
                TextStyle(
                    bold=row.emphasis or name == "variance",
                    color=colour,
                    align=PP_ALIGN.RIGHT if name in NUMERIC else PP_ALIGN.LEFT,
                    anchor=MSO_ANCHOR.MIDDLE,
                    line_spacing=0.95,
                    space_after_pt=0.0,
                ),
                f"financials/{row.category}/{name}",
            )
        )


def render(prs: PresentationType, spec: FinancialSummarySpec, page: int) -> Slide:
    slide = add_slide(prs)

    has_panel = bool(spec.considerations)
    total_w = T.body_width_with_panel() if has_panel else T.content_width()
    title_block(slide, spec.title, spec.subtitle, where=f"page {page}", width=total_w)

    x, top = T.content_left(), T.content_top()
    columns = _columns(x, total_w)

    add_rect(slide, x, top, total_w, T.FIN_HEADER_H, fill=T.PRIMARY, name="fin:header")
    headings = []
    for key, (cx, cw) in columns.items():
        box = add_textbox(
            slide, cx + T.FIN_PAD, top, cw - 2 * T.FIN_PAD, T.FIN_HEADER_H,
            name=f"fin:head-{key}",
        )
        headings.append(
            FitRequest(
                box.text_frame, [HEADINGS[key]], cw - 2 * T.FIN_PAD, T.FIN_HEADER_H,
                TextStyle(
                    bold=True, color=T.WHITE,
                    align=PP_ALIGN.RIGHT if key in NUMERIC else PP_ALIGN.LEFT,
                    anchor=MSO_ANCHOR.MIDDLE, line_spacing=1.0, space_after_pt=0.0, wrap=False,
                ),
                f"financials/{HEADINGS[key]} heading",
            )
        )
    fit_group(headings, T.FS_MICRO - 2, T.FS_MICRO)

    body = [*spec.rows, spec.total]
    gauge_block = (
        len(spec.gauges) * (T.FIN_GAUGE_H + T.FIN_GAUGE_GAP) + T.FIN_GAUGE_GAP
        if spec.gauges
        else 0
    )
    available = T.content_height() - T.FIN_HEADER_H - gauge_block
    row_h = min(T.FIN_ROW_MAX_H, available // len(body))

    widest = max((abs(r.variance_value) for r in body), default=0.0)
    bar_half = columns["bar"][1] // 2 - T.FIN_PAD
    scale = (bar_half / widest) if widest else 0.0

    rows_top = top + T.FIN_HEADER_H
    requests: dict[str, list[FitRequest]] = {"category": [], "figure": []}
    for i, row in enumerate(body):
        y = rows_top + i * row_h
        if row.emphasis:
            add_rect(slide, x, y, total_w, row_h, fill=T.tint(T.PRIMARY, 0.90),
                     name=f"surface:fin{i}")
        elif i % 2 == 0:
            add_rect(slide, x, y, total_w, row_h, fill=T.ROW_TINT, name=f"surface:fin{i}")
        _row(slide, row, columns, y, row_h, scale, f"fin{i}", requests)

    fit_group(requests["category"], T.FS_MIN_DENSE, T.FS_BODY)
    fit_group(requests["figure"], T.FS_MIN_DENSE, T.FS_DENSE)

    if spec.gauges:
        gauge_top = rows_top + row_h * len(body) + T.FIN_GAUGE_GAP
        fit_group(
            _gauges(slide, spec.gauges, x, gauge_top, total_w), T.FS_MICRO - 1, T.FS_DENSE
        )

    if has_panel:
        considerations_panel(slide, spec.considerations, where=f"page {page}/considerations")

    footer(slide, page)
    return slide
