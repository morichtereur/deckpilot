"""One card per work package: rating, progress, what is happening, what is next.

The progress bar is two stacked rectangles rather than a chart, so it stays a
shape a reader can drag - and so it renders identically everywhere, which a
charting object does not.
"""

from __future__ import annotations

from pptx.enum.text import MSO_ANCHOR, PP_ALIGN
from pptx.presentation import Presentation as PresentationType
from pptx.slide import Slide
from pptx.util import Emu

from deckpilot.renderer.base import (
    FitRequest,
    TextStyle,
    add_rect,
    add_slide,
    add_textbox,
    considerations_panel,
    fit_group,
    fit_text,
    footer,
    number_badge,
    title_block,
)
from deckpilot.specgen.schema import StatusCard, StatusOverviewSpec
from deckpilot.theme import tokens as T

ACTIVITY_STYLE = TextStyle(
    color=T.GRAY_DARK,
    bullet="•",
    bullet_indent=T.BULLET_INDENT,
    anchor=MSO_ANCHOR.TOP,
    line_spacing=T.LINE_SPACING,
    space_after_pt=T.SPACE_AFTER_PT,
)


def _grid(count: int) -> tuple[int, int]:
    """Cards read best two-up; three columns only once there are more than four."""
    columns = 2 if count <= 4 else 3
    rows = -(-count // columns)
    return columns, rows


def _progress_bar(slide: Slide, x: int, y: int, w: int, pct: int, key: str) -> None:
    """Two stacked rectangles: the track, then the filled portion."""
    label_w = T.STATUS_BAR_LABEL_W
    track_w = w - label_w - T.TIGHT_GAP
    add_rect(slide, x, y, track_w, T.STATUS_BAR_H, fill=T.GRAY_LIGHT, name=f"{key}:track")
    filled = int(track_w * pct / 100)
    if filled > 0:
        add_rect(slide, x, y, filled, T.STATUS_BAR_H, fill=T.SECONDARY, name=f"marker:{key}fill")

    box = add_textbox(
        slide, x + track_w + T.TIGHT_GAP, y, label_w, T.STATUS_BAR_H, name=f"{key}:pct",
    )
    fit_text(
        box.text_frame,
        f"{pct}%",
        T.FS_MICRO - 1,
        T.FS_DENSE,
        where=f"status/{key}/progress",
        style=TextStyle(
            bold=True, color=T.PRIMARY, align=PP_ALIGN.RIGHT, anchor=MSO_ANCHOR.MIDDLE,
            line_spacing=1.0, space_after_pt=0.0, wrap=False,
        ),
    )
    box.text_frame.margin_left = box.text_frame.margin_right = Emu(0)


def _card(slide: Slide, card: StatusCard, x: int, y: int, w: int, h: int,
          key: str) -> tuple[FitRequest, FitRequest, FitRequest | None]:
    add_rect(slide, x, y, w, h, fill=T.WHITE, line=T.GRAY_LIGHT, name=f"{key}:box")
    add_rect(slide, x, y, w, T.STATUS_CARD_HEAD_H, fill=T.PRIMARY, name=f"{key}:head")

    badge_y = y + (T.STATUS_CARD_HEAD_H - T.BADGE_D) // 2
    number_badge(slide, x + T.STATUS_PAD, badge_y, card.number,
                 fill=T.WHITE, text_color=T.PRIMARY, name=f"{key}:badge")

    rag_colour = T.RAG_COLORS[card.rag]
    chip_x = x + w - T.STATUS_PAD - T.CHIP_W
    from deckpilot.renderer.base import status_chip

    status_chip(
        slide, chip_x, y + (T.STATUS_CARD_HEAD_H - T.CHIP_H) // 2,
        card.rag.upper(), rag_colour, name=f"{key}:rag",
    )

    name_x = x + T.STATUS_PAD + T.BADGE_D + T.LABEL_GAP
    name_w = chip_x - name_x - T.TIGHT_GAP
    name = add_textbox(slide, name_x, y, name_w, T.STATUS_CARD_HEAD_H, name=f"{key}:name")
    header = FitRequest(
        name.text_frame, [card.name], name_w, T.STATUS_CARD_HEAD_H,
        TextStyle(bold=True, color=T.WHITE, anchor=MSO_ANCHOR.MIDDLE,
                  line_spacing=0.95, space_after_pt=0.0),
        f"status/{card.name}",
    )

    inner_x = x + T.STATUS_PAD
    inner_w = w - 2 * T.STATUS_PAD
    bar_y = y + T.STATUS_CARD_HEAD_H + T.STATUS_PAD
    _progress_bar(slide, inner_x, bar_y, inner_w, card.progress_pct, key)

    milestone_h = T.STATUS_MILESTONE_H if card.next_milestone else 0
    activities_y = bar_y + T.STATUS_BAR_H + T.STATUS_SECTION_GAP
    activities_h = (y + h) - activities_y - T.STATUS_PAD - milestone_h
    cell = add_textbox(slide, inner_x, activities_y, inner_w, activities_h,
                       name=f"{key}:activities")
    activities = FitRequest(
        cell.text_frame, list(card.activities), inner_w, activities_h,
        ACTIVITY_STYLE, f"status/{card.name}/activities",
    )

    milestone = None
    if card.next_milestone:
        box = add_textbox(
            slide, inner_x, activities_y + activities_h, inner_w, milestone_h,
            name=f"{key}:next",
        )
        milestone = FitRequest(
            box.text_frame, [f"Next: {card.next_milestone}"], inner_w, milestone_h,
            TextStyle(bold=True, color=T.SECONDARY, anchor=MSO_ANCHOR.MIDDLE,
                      line_spacing=0.95, space_after_pt=0.0),
            f"status/{card.name}/next milestone",
        )

    return header, activities, milestone


def render(prs: PresentationType, spec: StatusOverviewSpec, page: int) -> Slide:
    slide = add_slide(prs)

    has_panel = bool(spec.considerations)
    total_w = T.body_width_with_panel() if has_panel else T.content_width()
    title_block(slide, spec.title, spec.subtitle, where=f"page {page}", width=total_w)

    columns, rows = _grid(len(spec.cards))
    cw = T.col_w(columns, total=total_w)
    ch = T.row_h(rows)

    headers, activities, milestones = [], [], []
    for i, card in enumerate(spec.cards):
        x = T.col_x(i % columns, columns, total=total_w)
        y = T.row_y(i // columns, rows)
        header, activity, milestone = _card(slide, card, x, y, cw, ch, f"card{i}")
        headers.append(header)
        activities.append(activity)
        if milestone:
            milestones.append(milestone)

    fit_group(headers, T.FS_MICRO, T.FS_COLUMN_HEADER)
    fit_group(activities, T.FS_MIN_DENSE, T.FS_BODY, fill=True)
    if milestones:
        fit_group(milestones, T.FS_MICRO - 1, T.FS_DENSE)

    if has_panel:
        considerations_panel(slide, spec.considerations, where=f"page {page}/considerations")

    footer(slide, page)
    return slide
