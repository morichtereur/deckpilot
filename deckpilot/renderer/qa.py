"""Geometry linter for a built deck.

Visual inspection catches what only eyes catch - crowding, awkward wrapping,
a colour that fights the page. It is bad at the mechanical faults, which is
exactly what this checks: shapes off the slide, shapes across the outer margin,
shapes colliding, neighbours too close.

Shapes whose name starts with `bleed:` are deliberate full-bleed elements and
are exempt from the margin rule.
"""

from __future__ import annotations

from dataclasses import dataclass
from pathlib import Path

from pptx import Presentation

from deckpilot.theme import tokens as T

BLEED_PREFIX = "bleed:"


@dataclass(frozen=True)
class Finding:
    slide: int
    severity: str  # "error" | "warning"
    rule: str
    detail: str

    def __str__(self) -> str:
        return f"  slide {self.slide:>2}  [{self.severity:<7}] {self.rule}: {self.detail}"


@dataclass(frozen=True)
class Rect:
    name: str
    left: int
    top: int
    width: int
    height: int

    @property
    def right(self) -> int:
        return self.left + self.width

    @property
    def bottom(self) -> int:
        return self.top + self.height

    @property
    def area(self) -> int:
        return max(0, self.width) * max(0, self.height)

    def intersection_area(self, other: Rect) -> int:
        dx = min(self.right, other.right) - max(self.left, other.left)
        dy = min(self.bottom, other.bottom) - max(self.top, other.top)
        return dx * dy if dx > 0 and dy > 0 else 0

    def contains(self, other: Rect, slack: int = 0) -> bool:
        return (
            self.left - slack <= other.left
            and self.top - slack <= other.top
            and other.right <= self.right + slack
            and other.bottom <= self.bottom + slack
        )


def _rects(slide) -> list[Rect]:
    out = []
    for shape in slide.shapes:
        if shape.width is None or shape.height is None:
            continue
        out.append(Rect(shape.name, shape.left, shape.top, shape.width, shape.height))
    return out


def _is_bleed(rect: Rect) -> bool:
    return rect.name.startswith(BLEED_PREFIX)


def _has_extent(rect: Rect) -> bool:
    return rect.width > 0 and rect.height > 0


def _family(rect: Rect) -> str | None:
    """Shapes named `family:part` belong to one composed element.

    A title and its subtitle are deliberately set tight against each other; that
    is typography, not crowding, so the gap rule does not apply within a family.
    """
    return rect.name.split(":", 1)[0] if ":" in rect.name else None


def check_slide(index: int, slide, *, overlap_tolerance: float = 0.12) -> list[Finding]:
    findings: list[Finding] = []
    rects = [r for r in _rects(slide) if _has_extent(r)]

    for r in rects:
        if _is_bleed(r):
            continue
        if r.left < 0 or r.top < 0 or r.right > T.SLIDE_W or r.bottom > T.SLIDE_H:
            findings.append(
                Finding(index, "error", "off-slide", f"{r.name} extends past the slide edge")
            )
            continue
        over = []
        if r.left < T.MARGIN:
            over.append(f"left by {(T.MARGIN - r.left) / T.EMU_PER_INCH:.2f}\"")
        if r.top < T.MARGIN:
            over.append(f"top by {(T.MARGIN - r.top) / T.EMU_PER_INCH:.2f}\"")
        if r.right > T.SLIDE_W - T.MARGIN:
            over.append(f"right by {(r.right - (T.SLIDE_W - T.MARGIN)) / T.EMU_PER_INCH:.2f}\"")
        if r.bottom > T.SLIDE_H - T.MARGIN:
            over.append(f"bottom by {(r.bottom - (T.SLIDE_H - T.MARGIN)) / T.EMU_PER_INCH:.2f}\"")
        if over:
            findings.append(
                Finding(
                    index,
                    "error",
                    "margin",
                    f"{r.name} crosses the outer margin: {', '.join(over)}",
                )
            )

    # Collisions. Full containment is deliberate (a label inside its box); a
    # partial overlap is not.
    for i, a in enumerate(rects):
        for b in rects[i + 1 :]:
            if _is_bleed(a) or _is_bleed(b):
                continue
            overlap = a.intersection_area(b)
            if overlap == 0:
                continue
            if a.contains(b) or b.contains(a):
                continue
            smaller = min(a.area, b.area)
            if smaller and overlap / smaller > overlap_tolerance:
                findings.append(
                    Finding(
                        index,
                        "error",
                        "overlap",
                        f"{a.name} and {b.name} overlap by "
                        f"{overlap / smaller:.0%} of the smaller shape",
                    )
                )

    # Crowding: shapes that nearly touch without touching read as a mistake.
    for i, a in enumerate(rects):
        for b in rects[i + 1 :]:
            if _is_bleed(a) or _is_bleed(b):
                continue
            if a.intersection_area(b):
                continue
            family = _family(a)
            if family is not None and family == _family(b):
                continue
            vertical_overlap = min(a.bottom, b.bottom) - max(a.top, b.top) > 0
            horizontal_overlap = min(a.right, b.right) - max(a.left, b.left) > 0
            gap = None
            if vertical_overlap:
                gap = max(a.left, b.left) - min(a.right, b.right)
            elif horizontal_overlap:
                gap = max(a.top, b.top) - min(a.bottom, b.bottom)
            if gap is not None and 0 < gap < T.MIN_SHAPE_GAP:
                findings.append(
                    Finding(
                        index,
                        "warning",
                        "crowding",
                        f"{a.name} and {b.name} are only "
                        f"{gap / T.EMU_PER_INCH:.3f}\" apart",
                    )
                )
    return findings


def check_deck(path: str | Path) -> list[Finding]:
    prs = Presentation(str(path))
    findings: list[Finding] = []
    for i, slide in enumerate(prs.slides, start=1):
        findings.extend(check_slide(i, slide))
    return findings


def report(findings: list[Finding]) -> str:
    if not findings:
        return "Geometry check: clean."
    errors = [f for f in findings if f.severity == "error"]
    warnings = [f for f in findings if f.severity == "warning"]
    lines = [f"Geometry check: {len(errors)} error(s), {len(warnings)} warning(s)"]
    lines += [str(f) for f in errors + warnings]
    return "\n".join(lines)
