"""N columns, each a question and the characteristics that answer it.

The same shape covers stage gate criteria, a WHY / WHAT / HOW framing, or an
options comparison: a heading that poses something, and a list that resolves it.
The optional state tints the heading, which is what turns a set of criteria into
a status view without adding a second element to read.
"""

from __future__ import annotations

from pptx.enum.text import MSO_ANCHOR
from pptx.presentation import Presentation as PresentationType
from pptx.slide import Slide

from deckpilot.renderer.base import (
    FitRequest,
    TextStyle,
    add_rect,
    add_slide,
    add_textbox,
    considerations_panel,
    fit_group,
    fit_text,
    footer,
    title_block,
)
from deckpilot.specgen.schema import CriteriaColumnsSpec
from deckpilot.theme import tokens as T

STATE_COLORS = {
    "passed": T.STATUS_GREEN,
    "upcoming": T.SECONDARY,
    "at-risk": T.STATUS_RED,
    "neutral": T.PRIMARY,
}

BODY_STYLE = TextStyle(
    color=T.GRAY_DARK,
    bullet="•",
    bullet_indent=T.BULLET_INDENT,
    anchor=MSO_ANCHOR.TOP,
    line_spacing=T.LINE_SPACING,
    space_after_pt=T.SPACE_AFTER_PT,
)


def render(prs: PresentationType, spec: CriteriaColumnsSpec, page: int) -> Slide:
    slide = add_slide(prs)

    has_panel = bool(spec.considerations)
    total_w = T.body_width_with_panel() if has_panel else T.content_width()
    title_block(slide, spec.title, spec.subtitle, where=f"page {page}", width=total_w)

    n = len(spec.columns)
    cw = T.col_w(n, total=total_w)
    top, height = T.content_top(), T.content_height()

    questions, captions, bodies = [], [], []
    for i, column in enumerate(spec.columns):
        x = T.col_x(i, n, total=total_w)
        colour = STATE_COLORS[column.state]

        add_rect(slide, x, top, cw, height, fill=T.WHITE, line=T.GRAY_LIGHT, name=f"col{i}:box")
        add_rect(slide, x, top, cw, T.CRITERIA_HEAD_H, fill=colour, name=f"col{i}:head")
        inner_x = x + T.CRITERIA_PAD
        inner_w = cw - 2 * T.CRITERIA_PAD

        caption_h = T.CRITERIA_CAPTION_H if column.caption else 0
        question_h = T.CRITERIA_HEAD_H - caption_h
        question = add_textbox(slide, inner_x, top, inner_w, question_h, name=f"col{i}:question")
        questions.append(
            FitRequest(
                question.text_frame, [column.question], inner_w, question_h,
                TextStyle(
                    bold=True, color=T.on_color(colour),
                    anchor=MSO_ANCHOR.BOTTOM if column.caption else MSO_ANCHOR.MIDDLE,
                    line_spacing=0.95, space_after_pt=0.0,
                ),
                f"criteria/{column.question}",
            )
        )
        if column.caption:
            caption = add_textbox(
                slide, inner_x, top + question_h, inner_w, caption_h, name=f"col{i}:caption"
            )
            captions.append(
                FitRequest(
                    caption.text_frame, [column.caption], inner_w, caption_h,
                    TextStyle(
                        color=T.tint(T.on_color(colour), 0.25), anchor=MSO_ANCHOR.TOP,
                        line_spacing=1.0, space_after_pt=0.0, wrap=False,
                    ),
                    f"criteria/{column.question}/caption",
                )
            )

        label_y = top + T.CRITERIA_HEAD_H + T.CRITERIA_PAD
        label = add_textbox(slide, inner_x, label_y, inner_w, T.CRITERIA_LABEL_H,
                            name=f"col{i}:label")
        fit_text(
            label.text_frame,
            f"{spec.characteristics_label}:",
            T.FS_MICRO - 1,
            T.FS_DENSE,
            where=f"criteria/{column.question}/label",
            style=TextStyle(
                bold=True, color=T.SECONDARY, anchor=MSO_ANCHOR.MIDDLE,
                line_spacing=1.0, space_after_pt=0.0, wrap=False,
            ),
        )

        body_y = label_y + T.CRITERIA_LABEL_H
        body_h = top + height - body_y - T.CRITERIA_PAD
        body = add_textbox(slide, inner_x, body_y, inner_w, body_h, name=f"col{i}:body")
        bodies.append(
            FitRequest(
                body.text_frame, list(column.characteristics), inner_w, body_h,
                BODY_STYLE, f"criteria/{column.question}/characteristics",
            )
        )

    fit_group(questions, T.FS_MICRO, T.FS_COLUMN_HEADER)
    if captions:
        fit_group(captions, T.FS_MICRO - 2, T.FS_MICRO)
    fit_group(bodies, T.FS_MIN_DENSE, T.FS_BODY, fill=True)

    if has_panel:
        considerations_panel(slide, spec.considerations, where=f"page {page}/considerations")

    footer(slide, page)
    return slide
