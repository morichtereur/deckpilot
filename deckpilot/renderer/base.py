"""Shared building blocks for every layout.

Two rules hold across the whole renderer:

1. Everything is a native PowerPoint shape. No layout ever rasterises content,
   so a reader can select the text, recolour a bar or move a box.
2. No layout contains a literal dimension. Positions come from
   `deckpilot.theme.tokens`, and text is sized by `fit_text`, never guessed.
"""

from __future__ import annotations

import logging
from dataclasses import dataclass, replace

from pptx import Presentation
from pptx.dml.color import RGBColor
from pptx.enum.shapes import MSO_SHAPE
from pptx.enum.text import MSO_ANCHOR, PP_ALIGN
from pptx.oxml.ns import qn
from pptx.presentation import Presentation as PresentationType
from pptx.slide import Slide
from pptx.util import Emu, Pt

from deckpilot.renderer import text_metrics as tm
from deckpilot.theme import tokens as T

log = logging.getLogger("deckpilot.renderer")

BLANK_LAYOUT = 6  # the built-in blank layout, so nothing is inherited from a master


# --------------------------------------------------------------------------
# Deck and slide construction
# --------------------------------------------------------------------------


def new_deck() -> PresentationType:
    """A 16:9 presentation with no inherited placeholders."""
    prs = Presentation()
    prs.slide_width = Emu(T.SLIDE_W)
    prs.slide_height = Emu(T.SLIDE_H)
    return prs


def add_slide(prs: PresentationType) -> Slide:
    return prs.slides.add_slide(prs.slide_layouts[BLANK_LAYOUT])


# --------------------------------------------------------------------------
# Primitive shapes
# --------------------------------------------------------------------------


def add_rect(
    slide: Slide,
    x: int,
    y: int,
    w: int,
    h: int,
    *,
    fill: RGBColor | None = None,
    line: RGBColor | None = None,
    line_pt: float = T.HAIRLINE_PT,
    rounded: bool = False,
    shape_type: MSO_SHAPE | None = None,
    name: str | None = None,
):
    """A rectangle. `rounded` uses PowerPoint's rounded rectangle with the corner
    radius scaled to the shape's short side, so small chips do not look inflated."""
    kind = shape_type or (MSO_SHAPE.ROUNDED_RECTANGLE if rounded else MSO_SHAPE.RECTANGLE)
    shape = slide.shapes.add_shape(kind, Emu(x), Emu(y), Emu(w), Emu(h))
    shape.shadow.inherit = False

    if fill is None:
        shape.fill.background()
    else:
        shape.fill.solid()
        shape.fill.fore_color.rgb = fill

    if line is None:
        shape.line.fill.background()
    else:
        shape.line.color.rgb = line
        shape.line.width = Pt(line_pt)

    if name:
        shape.name = name

    if kind is MSO_SHAPE.ROUNDED_RECTANGLE:
        # adj is the radius as a fraction of the short side.
        shape.adjustments[0] = T.CORNER_RADIUS_RATIO

    shape.text_frame.word_wrap = True
    _clear(shape.text_frame)
    return shape


def add_line(
    slide: Slide, x1: int, y1: int, x2: int, y2: int, *, color: RGBColor, width_pt: float
):
    connector = slide.shapes.add_connector(1, Emu(x1), Emu(y1), Emu(x2), Emu(y2))  # 1 = straight
    connector.line.color.rgb = color
    connector.line.width = Pt(width_pt)
    connector.shadow.inherit = False
    return connector


def add_textbox(slide: Slide, x: int, y: int, w: int, h: int, *, name: str | None = None):
    box = slide.shapes.add_textbox(Emu(x), Emu(y), Emu(w), Emu(h))
    if name:
        box.name = name
    box.text_frame.word_wrap = True
    _clear(box.text_frame)
    return box


# --------------------------------------------------------------------------
# Text frames
# --------------------------------------------------------------------------


def _clear(frame) -> None:
    frame.clear()
    frame.paragraphs[0].text = ""


def set_insets(frame, inset: int = T.TEXT_INSET) -> None:
    frame.margin_left = frame.margin_right = Emu(inset)
    frame.margin_top = frame.margin_bottom = Emu(inset)


def usable_width_pt(shape_w: int, inset: int = T.TEXT_INSET, indent: int = 0) -> float:
    return max(1.0, (shape_w - 2 * inset - indent) / T.EMU_PER_PT)


def usable_height_pt(shape_h: int, inset: int = T.TEXT_INSET) -> float:
    return max(1.0, (shape_h - 2 * inset) / T.EMU_PER_PT)


def set_vertical_text(frame, direction: str = "vert270") -> None:
    """Rotate a text frame's flow. `vert270` reads bottom-to-top, which is the
    convention for a side label running up the left edge of a band."""
    frame._txBody.bodyPr.set("vert", direction)


def _set_bullet(paragraph, char: str, indent_emu: int, color: RGBColor) -> None:
    """Give a paragraph a real PowerPoint bullet with a hanging indent."""
    pPr = paragraph._p.get_or_add_pPr()
    pPr.set("marL", str(indent_emu))
    pPr.set("indent", str(-indent_emu))
    for tag in ("a:buNone", "a:buChar", "a:buAutoNum"):
        for existing in pPr.findall(qn(tag)):
            pPr.remove(existing)
    bu_clr = pPr.makeelement(qn("a:buClr"), {})
    srgb = pPr.makeelement(qn("a:srgbClr"), {"val": str(color)})
    bu_clr.append(srgb)
    pPr.append(bu_clr)
    bu = pPr.makeelement(qn("a:buChar"), {"char": char})
    pPr.append(bu)


def _style_run(run, size_pt: float, bold: bool, color: RGBColor, font: str = T.FONT) -> None:
    run.font.name = font
    run.font.size = Pt(size_pt)
    run.font.bold = bold
    run.font.color.rgb = color


@dataclass(frozen=True)
class TextStyle:
    size_pt: float = T.FS_BODY
    min_pt: float = T.FS_MIN_BODY
    bold: bool = False
    color: RGBColor = T.GRAY_DARK
    align: PP_ALIGN = PP_ALIGN.LEFT
    anchor: MSO_ANCHOR = MSO_ANCHOR.TOP
    line_spacing: float = T.LINE_SPACING
    space_after_pt: float = T.SPACE_AFTER_PT
    bullet: str | None = None
    bullet_indent: int = T.inches(0.13)
    wrap: bool = True  # False keeps a label on one line and shrinks it instead


def _avail(box_w: int, box_h: int, style: TextStyle) -> tuple[float, float]:
    indent = style.bullet_indent if style.bullet else 0
    return usable_width_pt(box_w, indent=indent), usable_height_pt(box_h)


def natural_height_pt(
    paragraphs: list[str], box_w: int, size_pt: float, style: TextStyle
) -> float:
    """Height in points this content wants at `size_pt`, including frame insets."""
    avail_w, _ = _avail(box_w, 0, style)
    body = tm.block_height(
        paragraphs, avail_w, size_pt, style.bold, style.line_spacing, style.space_after_pt
    )
    return body + 2 * (T.TEXT_INSET / T.EMU_PER_PT)


def _fits(paragraphs: list[str], avail_w: float, avail_h: float, size: float,
          style: TextStyle) -> bool:
    if not style.wrap:
        # One line per paragraph, so width is the binding constraint.
        widest = max(tm.text_width(p, size, style.bold) for p in paragraphs)
        height = tm.block_height(
            [""] * len(paragraphs), avail_w, size, style.bold, style.line_spacing,
            style.space_after_pt,
        )
        return widest <= avail_w and height <= avail_h
    if any(tm.widest_word(p, size, style.bold) > avail_w for p in paragraphs):
        return False  # a word too wide for the box would overhang it
    return (
        tm.block_height(
            paragraphs, avail_w, size, style.bold, style.line_spacing, style.space_after_pt
        )
        <= avail_h
    )


def _largest_fitting(check, min_pt: float, max_pt: float) -> float:
    """Largest half-point size in [min_pt, max_pt] for which `check(size)` holds.

    Returns `min_pt` when nothing fits; the caller decides what to do about that.
    """
    if check(max_pt):
        return float(max_pt)
    lo, hi = float(min_pt), float(max_pt)
    if not check(lo):
        return lo
    while hi - lo > 0.5:
        mid = round(((lo + hi) / 2) * 2) / 2
        if mid <= lo or mid >= hi:
            break
        if check(mid):
            lo = mid
        else:
            hi = mid
    return lo


@dataclass
class FitRequest:
    """One text box taking part in a group fit."""

    frame: object
    paragraphs: list[str]
    width: int
    height: int
    style: TextStyle
    where: str


def _place(req: FitRequest, size: float) -> None:
    """Write one request at `size`, truncating and warning if it still does not fit."""
    avail_w, avail_h = _avail(req.width, req.height, req.style)
    paragraphs = req.paragraphs
    truncated = False
    if not _fits(paragraphs, avail_w, avail_h, size, req.style) and req.style.wrap:
        line_h = size * tm.LINE_HEIGHT_FACTOR * req.style.line_spacing
        budget = int(
            (avail_h - req.style.space_after_pt * (len(paragraphs) - 1)) // line_h
        )
        per = max(1, max(budget, len(paragraphs)) // len(paragraphs))
        shortened = [
            tm.truncate(para, avail_w, size, per, req.style.bold) for para in paragraphs
        ]
        truncated = shortened != paragraphs
        paragraphs = shortened
    _write(req.frame, paragraphs, size, req.style)
    if truncated:
        log.warning(
            "%s: content truncated at %.1fpt - text does not fit its box "
            "(%.0fx%.0f pt available)",
            req.where,
            size,
            avail_w,
            avail_h,
        )


MAX_EXTRA_SPACE_AFTER_PT = 10.0


def _distribute_slack(live: list[FitRequest], size: float) -> None:
    """Push leftover vertical space into the gaps between bullets.

    A list that stops a third of the way up its box looks like a box that was
    sized wrong. Opening the bullet gaps instead makes the same content fill the
    same box and read as a deliberate rhythm. The extra is shared across the
    group - one value for the whole band - so peer columns stay on the same
    baselines. It is set by the fullest cell, which has the least to give.
    """
    extras = []
    for r in live:
        gaps = len(r.paragraphs) - 1
        if gaps <= 0:
            continue
        avail_w, avail_h = _avail(r.width, r.height, r.style)
        used = tm.block_height(
            r.paragraphs, avail_w, size, r.style.bold, r.style.line_spacing,
            r.style.space_after_pt,
        )
        extras.append(max(0.0, (avail_h - used) / gaps))
    if not extras:
        return
    extra = min(min(extras), MAX_EXTRA_SPACE_AFTER_PT)
    if extra <= 0.25:
        return
    for i, r in enumerate(live):
        live[i] = replace(
            r, style=replace(r.style, space_after_pt=r.style.space_after_pt + extra)
        )


def fit_group(
    requests: list[FitRequest], min_pt: float, max_pt: float, *, fill: bool = False
) -> float:
    """Set every request to the largest single size at which all of them fit.

    Peer elements - the headers of a column set, the cells of one band - must
    share a type size. Fitting each independently produces a slide where one
    column's text is a point smaller than its neighbour's, which reads as a
    mistake even when nothing overflows.

    With `fill`, leftover height is then spread across the bullet gaps rather
    than left at the bottom of the box.
    """
    live = [r for r in requests if r.paragraphs]
    if not live:
        for req in requests:
            _clear(req.frame)
        return float(max_pt)

    def check(size: float) -> bool:
        return all(_fits(r.paragraphs, *_avail(r.width, r.height, r.style), size, r.style)
                   for r in live)

    size = _largest_fitting(check, min_pt, max_pt)
    if fill:
        _distribute_slack(live, size)
    placed = {id(r.frame): r for r in live}
    for req in requests:
        if req.paragraphs:
            _place(placed.get(id(req.frame), req), size)
        else:
            _clear(req.frame)
    return size


def fit_text(
    frame,
    text: str | list[str],
    min_pt: float,
    max_pt: float,
    *,
    where: str = "unnamed element",
    style: TextStyle | None = None,
    shape_w: int | None = None,
    shape_h: int | None = None,
) -> float:
    """Write `text` into `frame` at the largest size between `min_pt` and `max_pt`
    that fits, then return the size used.

    Binary search down to `min_pt`; if the content still does not fit at the
    minimum, truncate with an ellipsis and log a WARNING naming the slide and
    element. A shipped deck should produce no such warnings - they mean the
    content and the box disagree, which is a content or layout decision, not
    something to paper over by clipping.
    """
    style = style or TextStyle()
    paragraphs = [text] if isinstance(text, str) else [p for p in text if p is not None]
    if not paragraphs:
        _clear(frame)
        return float(max_pt)

    shape = frame._parent
    req = FitRequest(
        frame=frame,
        paragraphs=paragraphs,
        width=shape_w if shape_w is not None else shape.width,
        height=shape_h if shape_h is not None else shape.height,
        style=style,
        where=where,
    )
    avail_w, avail_h = _avail(req.width, req.height, style)
    size = _largest_fitting(
        lambda s: _fits(paragraphs, avail_w, avail_h, s, style), min_pt, max_pt
    )
    _place(req, size)
    return size


def _write(frame, paragraphs: list[str], size_pt: float, style: TextStyle) -> None:
    frame.word_wrap = style.wrap
    frame.vertical_anchor = style.anchor
    set_insets(frame)
    _clear(frame)
    for i, text in enumerate(paragraphs):
        para = frame.paragraphs[0] if i == 0 else frame.add_paragraph()
        para.alignment = style.align
        para.line_spacing = style.line_spacing
        if i < len(paragraphs) - 1:
            para.space_after = Pt(style.space_after_pt)
        run = para.add_run()
        run.text = text
        _style_run(run, size_pt, style.bold, style.color)
        if style.bullet:
            _set_bullet(para, style.bullet, style.bullet_indent, style.color)


# --------------------------------------------------------------------------
# Shared slide furniture
# --------------------------------------------------------------------------


def title_block(
    slide: Slide,
    title: str,
    subtitle: str | None = None,
    *,
    where: str = "title",
    width: int | None = None,
) -> None:
    """The action title plus a one-line grey subtitle.

    Deliberately no rule beneath it: the type does the separating.
    """
    width = width if width is not None else T.content_width()

    box = add_textbox(slide, T.content_left(), T.MARGIN, width, T.TITLE_H, name="title:action")
    fit_text(
        box.text_frame,
        title,
        T.FS_MIN_TITLE,
        T.FS_TITLE,
        where=f"{where}/action title",
        style=TextStyle(
            bold=True,
            color=T.PRIMARY,
            anchor=MSO_ANCHOR.BOTTOM,
            line_spacing=0.95,
            space_after_pt=0.0,
        ),
    )

    if subtitle:
        sub = add_textbox(
            slide,
            T.content_left(),
            T.MARGIN + T.TITLE_H + T.TITLE_SUBTITLE_GAP,
            width,
            T.SUBTITLE_H,
            name="title:subtitle",
        )
        fit_text(
            sub.text_frame,
            subtitle,
            T.FS_MICRO,
            T.FS_SUBTITLE,
            where=f"{where}/subtitle",
            style=TextStyle(
                color=T.SUBTITLE_GRAY,
                anchor=MSO_ANCHOR.TOP,
                line_spacing=1.0,
                space_after_pt=0.0,
            ),
        )


def footer(slide: Slide, page: int, note: str = T.FOOTER_NOTE) -> None:
    """Confidentiality note on the left, page number on the right."""
    half = T.content_width() // 2
    left = add_textbox(
        slide, T.content_left(), T.footer_top(), half, T.FOOTER_H, name="footer:note"
    )
    fit_text(
        left.text_frame,
        note,
        T.FS_MICRO - 1,
        T.FS_MICRO,
        where=f"page {page}/footer note",
        style=TextStyle(
            color=T.FOOTER_GRAY, anchor=MSO_ANCHOR.MIDDLE, line_spacing=1.0, space_after_pt=0.0
        ),
    )
    right = add_textbox(
        slide,
        T.content_left() + T.content_width() - half,
        T.footer_top(),
        half,
        T.FOOTER_H,
        name="footer:page",
    )
    fit_text(
        right.text_frame,
        str(page),
        T.FS_MICRO,
        T.FS_MICRO,
        where=f"page {page}/page number",
        style=TextStyle(
            color=T.FOOTER_GRAY,
            align=PP_ALIGN.RIGHT,
            anchor=MSO_ANCHOR.MIDDLE,
            line_spacing=1.0,
            space_after_pt=0.0,
        ),
    )


def number_badge(
    slide: Slide,
    x: int,
    y: int,
    label: str,
    *,
    diameter: int = T.BADGE_D,
    fill: RGBColor = T.SECONDARY,
    text_color: RGBColor = T.WHITE,
    name: str | None = None,
):
    badge = add_rect(
        slide,
        x,
        y,
        diameter,
        diameter,
        fill=fill,
        shape_type=MSO_SHAPE.OVAL,
        name=name or f"badge:{label}",
    )
    fit_text(
        badge.text_frame,
        label,
        T.FS_MICRO - 2,
        T.FS_MICRO + 1,
        where=f"badge {label}",
        style=TextStyle(
            bold=True,
            color=text_color,
            align=PP_ALIGN.CENTER,
            anchor=MSO_ANCHOR.MIDDLE,
            line_spacing=1.0,
            space_after_pt=0.0,
        ),
    )
    badge.text_frame.margin_left = badge.text_frame.margin_right = Emu(0)
    badge.text_frame.margin_top = badge.text_frame.margin_bottom = Emu(0)
    return badge


def status_chip(
    slide: Slide,
    x: int,
    y: int,
    label: str,
    color: RGBColor,
    *,
    w: int = T.CHIP_W,
    h: int = T.CHIP_H,
    text_color: RGBColor = T.WHITE,
    name: str | None = None,
):
    chip = add_rect(slide, x, y, w, h, fill=color, rounded=True, name=name or f"chip:{label}")
    fit_text(
        chip.text_frame,
        label,
        T.FS_MICRO - 2,
        T.FS_MICRO,
        where=f"chip {label}",
        style=TextStyle(
            bold=True,
            color=text_color,
            align=PP_ALIGN.CENTER,
            anchor=MSO_ANCHOR.MIDDLE,
            line_spacing=1.0,
            space_after_pt=0.0,
        ),
    )
    chip.text_frame.margin_left = chip.text_frame.margin_right = Emu(0)
    chip.text_frame.margin_top = chip.text_frame.margin_bottom = Emu(0)
    return chip


def considerations_panel(
    slide: Slide,
    items: list[str],
    *,
    heading: str = "Considerations",
    top: int | None = None,
    height: int | None = None,
    where: str = "considerations",
) -> None:
    """The optional right-hand annotation panel, attachable to any content layout.

    Callers narrow their own body area with `tokens.body_width_with_panel()`;
    this function owns everything from `panel_left()` rightwards.
    """
    if not items:
        return
    top = T.content_top() if top is None else top
    height = T.content_height() if height is None else height
    x, w = T.panel_left(), T.panel_width()

    add_rect(slide, x, top, w, height, fill=T.PANEL_BG, name="panel:bg")

    head_h = T.inches(0.24)
    head = add_textbox(
        slide, x + T.PANEL_PAD, top + T.PANEL_PAD, w - 2 * T.PANEL_PAD, head_h, name="panel:heading"
    )
    fit_text(
        head.text_frame,
        heading,
        T.FS_MICRO,
        T.FS_BODY,
        where=f"{where}/heading",
        style=TextStyle(
            bold=True,
            color=T.PRIMARY,
            anchor=MSO_ANCHOR.MIDDLE,
            line_spacing=1.0,
            space_after_pt=0.0,
        ),
    )

    body_y = top + T.PANEL_PAD + head_h + T.inches(0.06)
    body_h = height - (body_y - top) - T.PANEL_PAD
    body = add_textbox(
        slide, x + T.PANEL_PAD, body_y, w - 2 * T.PANEL_PAD, body_h, name="panel:body"
    )
    fit_text(
        body.text_frame,
        items,
        T.FS_MIN_DENSE,
        T.FS_DENSE,
        where=f"{where}/body",
        style=TextStyle(
            color=T.GRAY_DARK,
            bullet="•",
            line_spacing=T.LINE_SPACING,
            space_after_pt=T.SPACE_AFTER_PT + 1,
        ),
    )
