"""Contents: the sections, and where each one starts.

The page numbers are the point. They can only be settled once the whole deck has
been assembled - which, because the appendix is paginated by measurement, means
after the renderer has been asked how many slides the RAID log needs. A contents
page with numbers that do not match the deck is worse than no contents page.
"""

from __future__ import annotations

from pptx.enum.text import MSO_ANCHOR, PP_ALIGN
from pptx.presentation import Presentation as PresentationType
from pptx.slide import Slide
from pptx.util import Emu

from deckpilot.renderer.base import (
    FitRequest,
    TextStyle,
    add_rect,
    add_slide,
    add_textbox,
    fit_group,
    footer,
    title_block,
)
from deckpilot.specgen.schema import AgendaSpec
from deckpilot.theme import tokens as T


def render(prs: PresentationType, spec: AgendaSpec, page: int) -> Slide:
    slide = add_slide(prs)
    title_block(slide, spec.title, spec.subtitle, where=f"page {page}")

    x, w = T.content_left(), T.content_width()
    top = T.content_top()
    row_h = min(T.AGENDA_ROW_MAX_H, T.content_height() // len(spec.entries))

    numbers, titles, captions, pages = [], [], [], []
    for i, entry in enumerate(spec.entries):
        y = top + i * row_h
        if i % 2 == 0:
            add_rect(slide, x, y, w, row_h, fill=T.ROW_TINT, name=f"surface:agenda{i}")

        number = add_textbox(
            slide, x + T.AGENDA_PAD, y, T.AGENDA_NUMBER_W, row_h, name=f"entry{i}:number"
        )
        numbers.append(
            FitRequest(
                number.text_frame, [entry.number], T.AGENDA_NUMBER_W, row_h,
                TextStyle(
                    bold=True, color=T.tint(T.PRIMARY, 0.55), anchor=MSO_ANCHOR.MIDDLE,
                    line_spacing=1.0, space_after_pt=0.0, wrap=False,
                ),
                f"agenda/{entry.title}/number",
            )
        )

        text_x = x + T.AGENDA_PAD + T.AGENDA_NUMBER_W
        text_w = w - (text_x - x) - T.AGENDA_PAGE_W - 2 * T.AGENDA_PAD
        caption_h = T.inches(0.20) if entry.caption else 0
        title_h = min(T.AGENDA_TITLE_H, row_h - caption_h)
        title_y = y + (row_h - title_h - caption_h) // 2

        heading = add_textbox(slide, text_x, title_y, text_w, title_h, name=f"entry{i}:title")
        titles.append(
            FitRequest(
                heading.text_frame, [entry.title], text_w, title_h,
                TextStyle(
                    bold=True, color=T.PRIMARY, anchor=MSO_ANCHOR.MIDDLE,
                    line_spacing=1.0, space_after_pt=0.0,
                ),
                f"agenda/{entry.title}",
            )
        )
        if entry.caption:
            caption = add_textbox(
                slide, text_x, title_y + title_h, text_w, caption_h, name=f"entry{i}:caption"
            )
            captions.append(
                FitRequest(
                    caption.text_frame, [entry.caption], text_w, caption_h,
                    TextStyle(
                        color=T.SUBTITLE_GRAY, anchor=MSO_ANCHOR.TOP,
                        line_spacing=1.0, space_after_pt=0.0,
                    ),
                    f"agenda/{entry.title}/caption",
                )
            )

        number_box = add_textbox(
            slide, x + w - T.AGENDA_PAGE_W - T.AGENDA_PAD, y, T.AGENDA_PAGE_W, row_h,
            name=f"entry{i}:page",
        )
        pages.append(
            FitRequest(
                number_box.text_frame, [str(entry.page)], T.AGENDA_PAGE_W, row_h,
                TextStyle(
                    bold=True, color=T.PRIMARY, align=PP_ALIGN.RIGHT,
                    anchor=MSO_ANCHOR.MIDDLE, line_spacing=1.0, space_after_pt=0.0, wrap=False,
                ),
                f"agenda/{entry.title}/page",
            )
        )
        number_box.text_frame.margin_right = Emu(0)

    fit_group(numbers, T.FS_SUBTITLE, T.FS_SECTION_TITLE)
    fit_group(titles, T.FS_BODY, T.FS_SUBTITLE + 2)
    if captions:
        fit_group(captions, T.FS_MICRO, T.FS_BODY)
    fit_group(pages, T.FS_BODY, T.FS_SUBTITLE + 2)

    footer(slide, page)
    return slide
