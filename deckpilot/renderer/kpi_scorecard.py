"""Benefit realisation: where each measure started, where it is, where it must get to.

The bar is the interesting part. A progress bar on its own says nothing - 29% of
the way to target is good or bad depending on how much of the work that produces
it has been done. So each bar carries a marker at the delivery progress of the
sub-stream responsible for the measure, and the bar is coloured by whether it has
reached that marker.

That is deliberate: benefits back-load. Judging them against elapsed calendar time
marks every measure late for the first two thirds of a programme, which tells the
reader nothing they can act on.
"""

from __future__ import annotations

from pptx.enum.text import MSO_ANCHOR, PP_ALIGN
from pptx.presentation import Presentation as PresentationType
from pptx.slide import Slide

from deckpilot.renderer.base import (
    FitRequest,
    TextStyle,
    add_rect,
    add_slide,
    add_textbox,
    considerations_panel,
    fit_group,
    fit_text,
    footer,
    title_block,
)
from deckpilot.specgen.schema import BenefitRow, KpiScorecardSpec
from deckpilot.theme import tokens as T

COLUMN_SHARE = {
    "name": 0.290,
    "owner": 0.150,
    "baseline": 0.062,
    "current": 0.062,
    "target": 0.062,
    "progress": 0.374,
}
HEADINGS = {
    "name": "Measure",
    "owner": "Owner",
    "baseline": "Baseline",
    "current": "Current",
    "target": "Target",
    "progress": "Progress from baseline to target",
}
NUMERIC = ("baseline", "current", "target")


def bar_colour(row: BenefitRow):
    """Green once the measure has caught its delivery, red when it is well adrift."""
    if row.attainment >= row.expected:
        return T.STATUS_GREEN
    if row.expected - row.attainment >= T.KPI_BEHIND_TOLERANCE:
        return T.STATUS_RED
    return T.STATUS_AMBER


def _columns(x: int, width: int) -> dict[str, tuple[int, int]]:
    out: dict[str, tuple[int, int]] = {}
    cursor = x
    for key, share in COLUMN_SHARE.items():
        w = int(width * share)
        out[key] = (cursor, w)
        cursor += w
    return out


def _progress(slide: Slide, row: BenefitRow, x: int, y: int, w: int, key: str) -> None:
    track_w = w - 2 * T.KPI_PAD
    track_x = x + T.KPI_PAD
    add_rect(slide, track_x, y, track_w, T.KPI_BAR_H, fill=T.GRAY_LIGHT, name=f"{key}:track")

    filled = int(track_w * row.attainment)
    if filled > 0:
        add_rect(
            slide, track_x, y, filled, T.KPI_BAR_H,
            fill=bar_colour(row), name=f"marker:{key}fill",
        )

    # The marker stands proud of the track so it reads as a threshold rather than
    # as part of the bar.
    tick_x = track_x + int(track_w * row.expected) - T.KPI_TICK_W // 2
    tick_x = min(max(tick_x, track_x), track_x + track_w - T.KPI_TICK_W)
    add_rect(
        slide,
        tick_x,
        y - T.KPI_TICK_OVERHANG,
        T.KPI_TICK_W,
        T.KPI_BAR_H + 2 * T.KPI_TICK_OVERHANG,
        fill=T.PRIMARY,
        name=f"marker:{key}expected",
    )


def render(prs: PresentationType, spec: KpiScorecardSpec, page: int) -> Slide:
    slide = add_slide(prs)

    has_panel = bool(spec.considerations)
    total_w = T.body_width_with_panel() if has_panel else T.content_width()
    title_block(slide, spec.title, spec.subtitle, where=f"page {page}", width=total_w)

    x, top = T.content_left(), T.content_top()
    columns = _columns(x, total_w)

    add_rect(slide, x, top, total_w, T.KPI_HEADER_H, fill=T.PRIMARY, name="kpi:header")
    heading_style = TextStyle(
        bold=True, color=T.WHITE, anchor=MSO_ANCHOR.MIDDLE,
        line_spacing=1.0, space_after_pt=0.0, wrap=False,
    )
    headings = []
    for key, (cx, cw) in columns.items():
        align = PP_ALIGN.RIGHT if key in NUMERIC else PP_ALIGN.LEFT
        box = add_textbox(
            slide, cx + T.KPI_PAD, top, cw - 2 * T.KPI_PAD, T.KPI_HEADER_H,
            name=f"kpi:head-{key}",
        )
        headings.append(
            FitRequest(
                box.text_frame, [HEADINGS[key]], cw - 2 * T.KPI_PAD, T.KPI_HEADER_H,
                TextStyle(**{**heading_style.__dict__, "align": align}),
                f"kpi/{HEADINGS[key]} heading",
            )
        )
    fit_group(headings, T.FS_MICRO - 2, T.FS_MICRO)

    rows_top = top + T.KPI_HEADER_H
    available = T.content_height() - T.KPI_HEADER_H
    row_h = min(T.KPI_ROW_MAX_H, available // len(spec.rows))

    names, owners, numbers = [], [], []
    for i, row in enumerate(spec.rows):
        y = rows_top + i * row_h
        if i % 2 == 0:
            add_rect(slide, x, y, total_w, row_h, fill=T.ROW_TINT, name=f"surface:kpi{i}")

        for key, (cx, cw) in columns.items():
            if key == "progress":
                _progress(slide, row, cx, y + (row_h - T.KPI_BAR_H) // 2, cw, f"kpi{i}")
                continue
            text = getattr(row, key)
            align = PP_ALIGN.RIGHT if key in NUMERIC else PP_ALIGN.LEFT
            bold = key == "current"
            box = add_textbox(
                slide, cx + T.KPI_PAD, y, cw - 2 * T.KPI_PAD, row_h, name=f"kpi{i}:{key}"
            )
            request = FitRequest(
                box.text_frame, [text], cw - 2 * T.KPI_PAD, row_h,
                TextStyle(
                    bold=bold,
                    color=T.PRIMARY if bold else T.GRAY_DARK,
                    align=align,
                    anchor=MSO_ANCHOR.MIDDLE,
                    line_spacing=0.95,
                    space_after_pt=0.0,
                ),
                f"kpi/{row.name}/{key}",
            )
            (names if key == "name" else owners if key == "owner" else numbers).append(request)

    # Names may wrap; owners and figures must not, so they are fitted separately.
    fit_group(names, T.FS_MIN_DENSE, T.FS_BODY)
    fit_group(owners, T.FS_MIN_DENSE, T.FS_DENSE)
    fit_group(numbers, T.FS_MIN_DENSE, T.FS_DENSE)

    # One line of key, because a marker nobody can decode is decoration.
    legend_y = rows_top + row_h * len(spec.rows) + T.KPI_LEGEND_GAP
    if legend_y + T.KPI_LEGEND_H <= T.content_bottom():
        legend = add_textbox(slide, x, legend_y, total_w, T.KPI_LEGEND_H, name="kpi:legend")
        fit_text(
            legend.text_frame,
            f"Marker shows {spec.expected_label.lower()}",
            T.FS_MICRO - 1,
            T.FS_MICRO,
            where="kpi/legend",
            style=TextStyle(
                color=T.FOOTER_GRAY, anchor=MSO_ANCHOR.MIDDLE,
                line_spacing=1.0, space_after_pt=0.0, wrap=False,
            ),
        )

    if has_panel:
        considerations_panel(slide, spec.considerations, where=f"page {page}/considerations")

    footer(slide, page)
    return slide
